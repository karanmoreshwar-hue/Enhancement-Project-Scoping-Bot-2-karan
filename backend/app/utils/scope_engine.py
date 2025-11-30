# app/utils/scope_engine.py
from __future__ import annotations
import json, re, logging, math, os, tempfile,anyio,pytesseract, openpyxl,tiktoken, pytz, graphviz,requests
from app import models
from calendar import monthrange
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
from io import BytesIO
from PIL import Image
from app.config.config import QDRANT_COLLECTION
from typing import Dict, Any, List
from datetime import datetime, timedelta
from app.utils import azure_blob
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.utils.ai_clients import (
    get_llm_client,
    get_embed_client,
    get_qdrant_client,
    embed_text_ollama,
)


logger = logging.getLogger(__name__)

# Init AI services
llm_cfg = get_llm_client()
embed_cfg = get_embed_client()
qdrant = get_qdrant_client()

# Utility function to round effort months to nearest 0.5
def round_to_half(value: float) -> float:
    """Round a number to the nearest 0.5 increment with minimum of 0.5.

    Examples:
        0.1 ‚Üí 0.5 (minimum)
        1.3 ‚Üí 1.5
        3.8 ‚Üí 4.0
        4.1 ‚Üí 4.0
        2.26 ‚Üí 2.5
    """
    rounded = round(value * 2) / 2
    return max(0.5, rounded)  # Ensure minimum of 0.5

def ollama_chat(prompt: str, model: str = llm_cfg["model"], temperature: float = 0.7, format_json: bool = False) -> str:
    """Call Ollama to generate text from a prompt.

    Args:
        prompt: The prompt to send to Ollama
        model: The model to use (default from config)
        temperature: Sampling temperature (default 0.7)
        format_json: If True, forces Ollama to return valid JSON (default False)
    """
    try:
        payload = {
            "model": model,
            "prompt": prompt,
            "temperature": temperature,
            "stream": False,
            "options": {
                "num_predict": -1,     # -1 means unlimited (let model decide)
                "num_ctx": 32768,      # Increase context window to 32K
                "temperature": temperature,
                "stop": []              # Remove any stop sequences that might truncate
            }
        }

        # Force JSON format output if requested
        if format_json:
            payload["format"] = "json"
            logger.info("üîß Ollama JSON format enforcement ENABLED (max tokens: UNLIMITED)")

        resp = requests.post(
            f"{llm_cfg['host']}/api/generate",
            json=payload,
            timeout=300  # Increase timeout to 5 minutes for longer responses
        )
        resp.raise_for_status()
        data = resp.json()
        return data.get("response", "").strip()
    except Exception as e:
        logger.error(f"Ollama chat failed: {e}")
        return ""


PROJECTS_BASE = "projects"


# Default Role Rates (USD/month)
ROLE_RATE_MAP: Dict[str, float] = {
    "Backend Developer": 3000.0,
    "Frontend Developer": 2800.0,
    "QA Analyst": 1800.0,
    "QA Engineer": 2000.0,
    "Data Engineer": 2800.0,
    "Data Analyst": 2200.0,
    "Data Architect": 3500.0,
    "UX Designer": 2500.0,
    "UI/UX Designer": 2600.0,
    "Project Manager": 3500.0,
    "Cloud Engineer": 3000.0,
    "BI Developer": 2700.0,
    "DevOps Engineer": 3200.0,
    "Security Administrator": 3000.0,
    "System Administrator": 2800.0,
    "Solution Architect": 4000.0,
}

#  helpers
def _strip_code_fences(s: str) -> str:
    m = re.search(r"```(?:json)?(.*?)```", s, flags=re.DOTALL | re.IGNORECASE)
    return m.group(1) if m else s

def _repair_json(text: str) -> str:
    """Attempt to fix common JSON syntax errors."""
    import re

    # Remove trailing commas before closing braces/brackets
    text = re.sub(r',\s*([}\]])', r'\1', text)

    # Fix missing commas between object elements (}{)
    text = re.sub(r'}\s*{', r'},{', text)

    # Fix missing commas between array elements (][)
    text = re.sub(r']\s*\[', r'],[', text)

    # Fix missing commas between object properties (common LLM error)
    # Match: "key": "value"<newline>"nextkey": where comma is missing
    text = re.sub(r'("\s*)\n\s*(")', r'\1,\n\2', text)

    # Fix unquoted keys (capture word followed by colon, add quotes)
    # Only match at start of line or after { or , to avoid false positives
    text = re.sub(r'([{,]\s*)([a-zA-Z_][a-zA-Z0-9_]*)\s*:', r'\1"\2":', text)

    return text


def _normalize_activity_fields(act: dict, activity_id: int) -> dict:
    """
    Normalize activity field names to match expected schema.
    Handles various field name variations from LLM.
    """
    from datetime import datetime, timedelta

    # Map various field name variations to expected field names
    activity_name = (
        act.get('Activities', '') or
        act.get('name', '') or
        act.get('activity', '') or
        act.get('Activity', '')
    )

    description = (
        act.get('Description', '') or
        act.get('description', '') or
        activity_name  # Fallback to activity name
    )

    owner = (
        act.get('Owner', '') or
        act.get('owner', '') or
        act.get('responsible', '') or
        act.get('assignee', '') or
        "Backend Developer"  # Default fallback
    )

    resources = act.get('Resources', '') or act.get('resources', '')
    if isinstance(resources, list):
        resources = ", ".join(resources)

    start_date = (
        act.get('Start Date', '') or
        act.get('start_date', '') or
        act.get('startDate', '')
    )

    end_date = (
        act.get('End Date', '') or
        act.get('end_date', '') or
        act.get('endDate', '')
    )

    # Handle effort/duration in various forms
    effort_months = (
        act.get('Effort Months', 0) or
        act.get('effort_months', 0) or
        act.get('effortMonths', 0) or
        act.get('duration', 0) or
        act.get('story_points', 0) / 20.0  # Convert story points to months (rough estimate)
    )

    if not effort_months or effort_months <= 0:
        effort_months = 1.0

    # Round effort months to nearest 0.5
    effort_months = round_to_half(float(effort_months))

    # If no dates provided, calculate from today
    if not start_date:
        start = datetime.today() + timedelta(days=(activity_id - 1) * 7)
        start_date = start.strftime("%Y-%m-%d")
        end_date = (start + timedelta(days=int(effort_months * 30))).strftime("%Y-%m-%d")

    return {
        "ID": activity_id,
        "Activities": activity_name,
        "Description": description,
        "Owner": owner,
        "Resources": resources,
        "Start Date": start_date,
        "End Date": end_date,
        "Effort Months": float(effort_months)
    }


def _transform_nested_to_flat_schema(raw: dict, project) -> dict:
    """
    Transform LLM's nested schema (with phases containing activities)
    into the flat schema expected by the backend.

    Handles cases where LLM returns:
    {
      "project": "...",
      "phases": [{..., "activities": [...]}],
      ...
    }

    And converts to:
    {
      "overview": {...},
      "activities": [...],
      "resourcing_plan": []
    }
    """
    # Check if already in flat format but need field normalization
    if raw.get('overview') and raw.get('activities'):
        logger.info("‚úÖ JSON has flat structure, checking field names...")
        # Normalize activity field names even if structure is correct
        activities = raw.get('activities', [])
        if activities and isinstance(activities, list):
            normalized_activities = []
            for idx, act in enumerate(activities, 1):
                if isinstance(act, dict):
                    normalized = _normalize_activity_fields(act, idx)
                    # Check if normalization was needed
                    if 'Activities' not in act or 'Owner' not in act:
                        logger.info(f"üîß Normalized activity {idx} field names")
                    normalized_activities.append(normalized)
                elif isinstance(act, str):
                    # Handle string activities
                    normalized_activities.append({
                        "ID": idx,
                        "Activities": act,
                        "Description": act,
                        "Owner": "Backend Developer",
                        "Resources": "",
                        "Start Date": "",
                        "End Date": "",
                        "Effort Months": 1.0
                    })

            if normalized_activities:
                raw['activities'] = normalized_activities
                logger.info(f"‚úÖ Normalized {len(normalized_activities)} activities")

        return raw

    # Check if data is wrapped in a "data" key - unwrap it
    if raw.get('data') and isinstance(raw.get('data'), dict):
        logger.info("üîì Unwrapping nested 'data' key...")
        raw = raw.get('data')

    # Check if it's in nested format (has 'phases' or 'project' at root level, or activities inside)
    if not (raw.get('phases') or raw.get('project') or raw.get('activities')):
        logger.warning("‚ö†Ô∏è JSON format unclear - returning as-is")
        return raw

    logger.info("üîÑ Transforming nested JSON schema to flat format...")

    # Helper function to safely extract string values from potentially nested dicts
    def safe_extract(key: str, default=''):
        """Extract value, handling both string and nested dict cases."""
        val = raw.get(key, default)
        if isinstance(val, dict):
            # If it's a dict, try to extract 'name' or the key itself
            return val.get('name', '') or val.get(key, default)
        return val or default

    # Build overview from top-level fields
    # Handle case where 'project' might be a dict with nested fields
    project_val = raw.get('project', '')
    if isinstance(project_val, dict):
        project_name = project_val.get('name', '') or getattr(project, 'name', '')
        logger.info(f"üìã Extracted project name from nested dict: {project_name}")
    else:
        project_name = project_val or getattr(project, 'name', '')

    overview = {
        "Project Name": project_name,
        "Domain": safe_extract('domain', getattr(project, 'domain', '')),
        "Complexity": safe_extract('complexity', getattr(project, 'complexity', '')),
        "Tech Stack": safe_extract('tech_stack', getattr(project, 'tech_stack', '')),
        "Use Cases": safe_extract('use_cases', getattr(project, 'use_cases', '')),
        "Compliance": safe_extract('compliance', getattr(project, 'compliance', '')),
        "Duration": raw.get('duration', 0) or getattr(project, 'duration', 0)
    }

    # Extract and flatten activities
    # Handle both:
    # 1. Activities nested in phases: {"phases": [{"activities": [...]}]}
    # 2. Activities directly in raw: {"activities": [...]}
    activities = []
    activity_id = 1

    # First check if activities are in phases
    phases = raw.get('phases', [])
    if phases:
        for phase in phases:
            phase_name = phase.get('name', 'Unnamed Phase')
            phase_activities = phase.get('activities', [])

            for act in phase_activities:
                # Handle both string and dict activities
                if isinstance(act, str):
                    # Activity is just a string description
                    flat_activity = {
                        "ID": activity_id,
                        "Activities": act,
                        "Description": act,  # Use same string for description
                        "Owner": "Backend Developer",
                        "Resources": "",
                        "Start Date": "",
                        "End Date": "",
                        "Effort Months": 1.0
                    }
                elif isinstance(act, dict):
                    # Activity is a dict with structured fields
                    flat_activity = {
                        "ID": activity_id,
                        "Activities": act.get('name', '') or act.get('activity', ''),
                        "Description": act.get('description', ''),
                        "Owner": act.get('owner', '') or act.get('responsible', '') or "Backend Developer",
                        "Resources": ", ".join(act.get('resources', [])) if isinstance(act.get('resources'), list) else act.get('resources', ''),
                        "Start Date": act.get('start_date', '') or act.get('startDate', ''),
                        "End Date": act.get('end_date', '') or act.get('endDate', ''),
                        "Effort Months": act.get('effort_months', 0) or act.get('effortMonths', 0) or 1.0
                    }
                else:
                    # Skip invalid activity types
                    logger.warning(f"‚ö†Ô∏è Skipping invalid activity type: {type(act)}")
                    continue

                # If no start/end dates, calculate from today
                if not flat_activity["Start Date"]:
                    from datetime import datetime, timedelta
                    start = datetime.today() + timedelta(days=(activity_id - 1) * 7)
                    flat_activity["Start Date"] = start.strftime("%Y-%m-%d")
                    flat_activity["End Date"] = (start + timedelta(days=30)).strftime("%Y-%m-%d")
                    flat_activity["Effort Months"] = 1.0

                activities.append(flat_activity)
                activity_id += 1

        logger.info(f"‚úÖ Extracted {len(activities)} activities from {len(phases)} phases")

    # If no phases, check if activities are directly at root level
    elif raw.get('activities'):
        logger.info("üìã Found activities directly at root level")
        for act in raw.get('activities', []):
            # Same handling as above
            if isinstance(act, str):
                flat_activity = {
                    "ID": activity_id,
                    "Activities": act,
                    "Description": act,
                    "Owner": "Backend Developer",
                    "Resources": "",
                    "Start Date": "",
                    "End Date": "",
                    "Effort Months": 1.0
                }
            elif isinstance(act, dict):
                flat_activity = {
                    "ID": activity_id,
                    "Activities": act.get('name', '') or act.get('activity', ''),
                    "Description": act.get('description', ''),
                    "Owner": act.get('owner', '') or act.get('responsible', '') or "Backend Developer",
                    "Resources": ", ".join(act.get('resources', [])) if isinstance(act.get('resources'), list) else act.get('resources', ''),
                    "Start Date": act.get('start_date', '') or act.get('startDate', ''),
                    "End Date": act.get('end_date', '') or act.get('endDate', ''),
                    "Effort Months": act.get('effort_months', 0) or act.get('effortMonths', 0) or 1.0
                }
            else:
                logger.warning(f"‚ö†Ô∏è Skipping invalid activity type: {type(act)}")
                continue

            # If no start/end dates, calculate from today
            if not flat_activity["Start Date"]:
                from datetime import datetime, timedelta
                start = datetime.today() + timedelta(days=(activity_id - 1) * 7)
                flat_activity["Start Date"] = start.strftime("%Y-%m-%d")
                flat_activity["End Date"] = (start + timedelta(days=30)).strftime("%Y-%m-%d")
                flat_activity["Effort Months"] = 1.0

            activities.append(flat_activity)
            activity_id += 1

        logger.info(f"‚úÖ Extracted {len(activities)} activities from root level")

    # Build transformed structure
    transformed = {
        "overview": overview,
        "activities": activities,
        "resourcing_plan": [],  # Will be auto-generated by clean_scope
    }

    # Preserve risks/assumptions if present
    if raw.get('risks'):
        transformed['risks'] = raw.get('risks')
    if raw.get('assumptions'):
        transformed['assumptions'] = raw.get('assumptions')
    if raw.get('project_summary'):
        transformed['project_summary'] = raw.get('project_summary')

    return transformed


def _extract_json(s: str) -> dict:
    raw = _strip_code_fences(s or "")
    try:
        parsed = json.loads(raw.strip())
        # If Ollama returns a list at root level, check if it's activities
        if isinstance(parsed, list):
            logger.warning(f"‚ö†Ô∏è  Ollama returned a list instead of dict. Wrapping in activities key.")
            return {"activities": parsed}
        return parsed if isinstance(parsed, dict) else {}
    except Exception as e:
        logger.warning(f"‚ö†Ô∏è  First JSON parse attempt failed: {str(e)}")
        logger.warning(f"   Trying to extract JSON from braces...")
        start, end = raw.find("{"), raw.rfind("}")
        if start >= 0 and end > start:
            try:
                extracted = raw[start:end+1]
                logger.info(f"   Extracted JSON length: {len(extracted)} chars")
                logger.info(f"   Extracted JSON preview (first 300 chars): {extracted[:300]}")
                logger.info(f"   Extracted JSON ending (last 200 chars): {extracted[-200:]}")
                parsed = json.loads(extracted)
                if isinstance(parsed, list):
                    logger.warning(f"‚ö†Ô∏è  Ollama returned a list instead of dict. Wrapping in activities key.")
                    return {"activities": parsed}
                logger.info(f"‚úÖ Successfully parsed JSON with {len(parsed)} top-level keys: {list(parsed.keys())}")
                return parsed if isinstance(parsed, dict) else {}
            except Exception as e2:
                logger.warning(f"‚ö†Ô∏è  Second JSON parse attempt also failed: {str(e2)}")
                logger.warning(f"   Attempting JSON repair...")
                try:
                    # Try to repair common JSON syntax errors
                    repaired = _repair_json(extracted)
                    logger.info(f"   Repaired JSON preview (first 300 chars): {repaired[:300]}")
                    logger.info(f"   Repaired JSON ending (last 200 chars): {repaired[-200:]}")
                    parsed = json.loads(repaired)
                    if isinstance(parsed, list):
                        logger.warning(f"‚ö†Ô∏è  Ollama returned a list instead of dict. Wrapping in activities key.")
                        return {"activities": parsed}
                    logger.info(f"‚úÖ Successfully parsed repaired JSON with {len(parsed)} top-level keys: {list(parsed.keys())}")
                    return parsed if isinstance(parsed, dict) else {}
                except Exception as e3:
                    logger.error(f"‚ùå JSON repair also failed: {str(e3)}")
                    logger.error(f"   Raw text length: {len(raw)} chars")
                    logger.error(f"   Raw text preview (first 300 chars): {raw[:300]}")
                    logger.error(f"   Raw text ending (last 200 chars): {raw[-200:]}")
                    return {}
        return {}



def _parse_date_safe(val: Any, fallback: datetime = None) -> datetime:
    """Try to parse a date string; return fallback if invalid."""
    if not val:
        return fallback
    try:
        return datetime.strptime(str(val), "%Y-%m-%d")
    except Exception:
        return fallback

def _safe_str(val: Any) -> str:
    """Convert value to string, handling arrays and dictionaries by joining them."""
    if val is None:
        return ""

    # If it's a dictionary, extract all values and flatten them
    if isinstance(val, dict):
        all_values = []
        for v in val.values():
            if isinstance(v, list):
                all_values.extend(v)
            elif v:
                all_values.append(str(v))
        return ", ".join(str(item).strip() for item in all_values if item)

    # If it's a list/array, join with commas
    if isinstance(val, list):
        return ", ".join(str(item).strip() for item in val if item)

    # Check if it's a string representation of an array like "['item1', 'item2']"
    if isinstance(val, str) and val.strip().startswith('[') and val.strip().endswith(']'):
        try:
            import json
            parsed = json.loads(val.replace("'", '"'))  # Convert single quotes to double quotes for JSON
            if isinstance(parsed, list):
                return ", ".join(str(item).strip() for item in parsed if item)
        except:
            # If JSON parsing fails, try Python literal eval
            try:
                import ast
                parsed = ast.literal_eval(val)
                if isinstance(parsed, list):
                    return ", ".join(str(item).strip() for item in parsed if item)
            except:
                pass  # If both fail, return as-is below

    # Check if it's a string representation of a dict like "{'key': ['val1', 'val2']}"
    if isinstance(val, str) and val.strip().startswith('{') and val.strip().endswith('}'):
        try:
            import ast
            parsed = ast.literal_eval(val)
            if isinstance(parsed, dict):
                all_values = []
                for v in parsed.values():
                    if isinstance(v, list):
                        all_values.extend(v)
                    elif v:
                        all_values.append(str(v))
                return ", ".join(str(item).strip() for item in all_values if item)
        except:
            pass  # If parsing fails, return as-is below

    return str(val).strip()

async def get_rate_map_for_project(db: AsyncSession, project) -> Dict[str, float]:
    """
    Fetch rate cards for the given project/company.
    Falls back to Sigmoid default rates if none exist
    """
    try:
        # If project has company_id, try fetching company-specific rate cards
        if getattr(project, "company_id", None):
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == project.company_id)
            )
            ratecards = result.scalars().all()
            if ratecards:
                return {r.role_name: float(r.monthly_rate) for r in ratecards}

        sigmoid_result = await db.execute(
            select(models.Company).filter(models.Company.name == "Sigmoid")
        )
        sigmoid = sigmoid_result.scalars().first()
        if sigmoid:
            result = await db.execute(
                select(models.RateCard)
                .filter(models.RateCard.company_id == sigmoid.id)
            )
            sigmoid_rates = result.scalars().all()
            if sigmoid_rates:
                return {r.role_name: float(r.monthly_rate) for r in sigmoid_rates}

    except Exception as e:
        logger.warning(f"Failed to fetch rate cards: {e}")
    return ROLE_RATE_MAP

def extract_text_from_file(file_bytes_io: BytesIO, file_name: str) -> str:
    """
    Extract text from a file given its bytes and filename.

    Args:
        file_bytes_io: BytesIO object containing file bytes
        file_name: Name of the file (used to determine file type)

    Returns:
        Extracted text content
    """
    suffix = os.path.splitext(file_name)[-1].lower()
    file_bytes = file_bytes_io.read()
    file_bytes_io.seek(0)  # Reset for potential re-reading

    content = ""
    try:
        if suffix == ".pdf":
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                content = extract_pdf_text(tmp_path)
            finally:
                os.remove(tmp_path)

        elif suffix == ".docx":
            doc = Document(BytesIO(file_bytes))
            content = "\n".join(p.text for p in doc.paragraphs)

        elif suffix == ".pptx":
            prs = Presentation(BytesIO(file_bytes))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = "\n".join(texts)

        elif suffix in [".xlsx", ".xlsm"]:
            wb = openpyxl.load_workbook(BytesIO(file_bytes))
            sheet = wb.active
            content = "\n".join(
                " ".join(str(cell) if cell else "" for cell in row)
                for row in sheet.iter_rows(values_only=True)
            )

        elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
            img = Image.open(BytesIO(file_bytes))
            content = pytesseract.image_to_string(img)

        else:
            # Try as text file
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name
            try:
                with open(tmp_path, "r", encoding="utf-8", errors="ignore") as fh:
                    content = fh.read()
            finally:
                os.remove(tmp_path)

    except Exception as e:
        logger.warning(f"Text extraction failed for {file_name}: {e}")

    return content.strip()


async def _extract_text_from_files(files: List[dict]) -> str:
    results: List[str] = []

    async def _extract_single(f: dict) -> None:
        try:
            blob_bytes = await azure_blob.download_bytes(f["file_path"])
            suffix = os.path.splitext(f["file_name"])[-1].lower()

            def process_file() -> str:
                content = ""
                try:
                    if suffix == ".pdf":
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            content = extract_pdf_text(tmp_path)
                        finally:
                            os.remove(tmp_path)

                    elif suffix == ".docx":
                        doc = Document(BytesIO(blob_bytes))
                        content = "\n".join(p.text for p in doc.paragraphs)

                    elif suffix == ".pptx":
                        prs = Presentation(BytesIO(blob_bytes))
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    texts.append(shape.text)
                        content = "\n".join(texts)

                    elif suffix in [".xlsx", ".xlsm"]:
                        wb = openpyxl.load_workbook(BytesIO(blob_bytes))
                        sheet = wb.active
                        content = "\n".join(
                            " ".join(str(cell) if cell else "" for cell in row)
                            for row in sheet.iter_rows(values_only=True)
                        )

                    elif suffix in [".png", ".jpg", ".jpeg", ".tiff"]:
                        img = Image.open(BytesIO(blob_bytes))
                        content = pytesseract.image_to_string(img)

                    else:
                        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                            tmp.write(blob_bytes)
                            tmp_path = tmp.name
                        try:
                            with open(tmp_path, "r", encoding="utf-8", errors="ignore") as fh:
                                content = fh.read()
                        finally:
                            os.remove(tmp_path)

                except Exception as e:
                    logger.warning(f"Extraction failed for {f['file_name']}: {e}")

                return content.strip()

            text = await anyio.to_thread.run_sync(process_file)

            if text:
                results.append(text)
            else:
                logger.warning(f"Extracted no text from {f['file_name']}")

        except Exception as e:
            logger.warning(f"Failed to extract {f.get('file_name')} (path={f.get('file_path')}): {e}")

    async with anyio.create_task_group() as tg:
        for f in files:
            tg.start_soon(_extract_single, f)

    return "\n\n".join(results)


def _rag_retrieve(query: str, k: int = 5) -> List[Dict]:
    """
    Retrieve semantically similar chunks from Qdrant for RAG.
    Uses Ollama embedding model and returns list of matched chunks.
    Skips retrieval if no valid embedding found.
    """
    try:
        q_emb_list = embed_text_ollama([query])

        # Skip if no valid embeddings returned
        if not q_emb_list or not q_emb_list[0]:
            logger.warning("‚ö†Ô∏è No valid embedding generated ‚Äî skipping Qdrant retrieval.")
            return []

        q_emb = q_emb_list[0]

        # Sanity check vector dimension
        if not isinstance(q_emb, list) or len(q_emb) == 0:
            logger.warning("‚ö†Ô∏è Empty embedding vector ‚Äî skipping retrieval.")
            return []

        client = get_qdrant_client()
        # Search ONLY in KB collection (case studies are in separate collection)
        results = client.search(
            collection_name=QDRANT_COLLECTION,  # KB documents only, no case studies
            query_vector=q_emb,
            limit=k,
            with_payload=True
        )

        logger.info(f"üîç Searching knowledge base (Qdrant) - found {len(results)} results (KB documents only, excluding case studies)")

        hits = []
        for r in results:
            payload = r.payload or {}
            file_name = payload.get("file_name", "unknown")
            chunk_index = payload.get("chunk_index", "?")
            score = r.score

            # Log each result with details
            logger.info(f"   üìÑ {file_name} (chunk {chunk_index}): similarity {score:.3f}")

            hits.append({
                "id": payload.get("chunk_id", str(r.id)),
                "parent_id": payload.get("parent_id"),
                "content": payload.get("chunk", ""),
                "title": payload.get("title", ""),
                "score": r.score,
            })

        # Group by parent_id for consistency
        grouped = {}
        for h in hits:
            grouped.setdefault(h["parent_id"], []).append({
                "id": h["id"],
                "content": h["content"],
                "title": h["title"],
                "score": h["score"],
            })

        return [
            {"parent_id": pid, "chunks": chs}
            for pid, chs in grouped.items()
        ]

    except Exception as e:
        logger.warning(f"RAG retrieval (Qdrant) failed: {e}")
        return []

def _build_scope_prompt(rfp_text: str, kb_chunks: List[str], project=None, questions_context: str | None = None, rate_card_roles: List[str] | None = None) -> str:
    import tiktoken

    # Tokenizer
    tokenizer = tiktoken.get_encoding("cl100k_base")
    # Safe token budget (128k, keep ~4k for completion & system messages)
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # Trim RFP text
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 3000:
        rfp_tokens = rfp_tokens[:3000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # Trim KB context
    safe_kb_chunks = []
    for ch in kb_chunks or []:
        tokens = tokenizer.encode(ch)
        if used_tokens + len(tokens) > max_total_tokens:
            break
        safe_kb_chunks.append(ch)
        used_tokens += len(tokens)

    kb_context = "\n\n".join(safe_kb_chunks) if safe_kb_chunks else "(no KB context found)"

    name = (getattr(project, "name", "") or "").strip()
    domain = (getattr(project, "domain", "") or "").strip()
    complexity = (getattr(project, "complexity", "") or "").strip()
    tech_stack = (getattr(project, "tech_stack", "") or "").strip()
    use_cases = (getattr(project, "use_cases", "") or "").strip()
    compliance = (getattr(project, "compliance", "") or "").strip()
    duration = str(getattr(project, "duration", "") or "").strip()

    user_context = (
        "Some overview fields have been provided by the user.\n"
        "Treat these user-provided values as the source of truth.\n"
        "Only fill in fields that are blank ‚Äî do NOT overwrite the given values.\n\n"
        f"Project Name: {name or '(infer if missing)'}\n"
        f"Domain: {domain or '(infer if missing)'}\n"
        f"Complexity: {complexity or '(infer if missing)'}\n"
        f"Tech Stack: {tech_stack or '(infer if missing)'}\n"
        f"Use Cases: {use_cases or '(infer if missing)'}\n"
        f"Compliance: {compliance or '(infer if missing)'}\n"
        f"Duration (months): {duration or '(infer if missing)'}\n\n"
    )

    today_str = datetime.today().date().isoformat()

    return (
        "CRITICAL INSTRUCTION: You MUST output ONLY valid JSON. Do NOT include any explanations, commentary, thinking process, or markdown.\n"
        "Do NOT start with 'Okay' or 'Here is' or any prose. Your ENTIRE response must be valid JSON and nothing else.\n\n"
        "You are an expert AI project planner.\n"
        "Use the RFP/project text as the **primary source** \n"
        "Use questions and answers to clarify ambiguities.\n"
        "but enrich missing fields with the Knowledge Base context (if relevant).\n\n"
        "‚ö†Ô∏è MANDATORY REQUIREMENT: You MUST generate a complete 'activities' array with at least 8-15 activities.\n"
        "‚ùå DO NOT generate empty activities array - this is UNACCEPTABLE.\n"
        "‚ùå DO NOT return ONLY metadata without activities - this is a CRITICAL ERROR.\n"
        "‚úÖ The 'activities' array is THE MOST IMPORTANT part of your response.\n\n"
        "Output schema (YOUR ENTIRE RESPONSE MUST MATCH THIS EXACT FORMAT):\n"
        "{\n"
        '  "overview": {\n'
        '    "Project Name": string,\n'
        '    "Domain": string,\n'
        '    "Complexity": string,\n'
        '    "Tech Stack": string,\n'
        '    "Use Cases": string,\n'
        '    "Compliance": string,\n'
        '    "Duration": number\n'
        "  },\n"
        '  "activities": [\n'
        '    {\n'
        '      "ID": int,\n'
        '      "Activities": string,\n'
        '      "Description": string | null,\n'
        '      "Owner": string | null,\n'
        '      "Resources": string | null,\n'
        '      "Start Date": "yyyy-mm-dd",\n'
        '      "End Date": "yyyy-mm-dd",\n'
        '      "Effort Months": number\n'
        "    }\n"
        "  ],\n"
        '  "resourcing_plan": [],\n'
        '  "project_summary": {\n'
        '    "executive_summary": string,\n'
        '    "key_deliverables": [string],\n'
        '    "success_criteria": [string],\n'
        '    "risks_and_mitigation": [{risk: string, mitigation: string}]\n'
        "  }\n"
        "}\n\n"
        "‚ö†Ô∏è CRITICAL: The 'activities' array MUST contain at least 8-15 detailed activities covering ALL project phases:\n"
        "   - Requirements gathering, analysis, and planning activities\n"
        "   - Design and architecture activities\n"
        "   - Development activities (broken down by feature/module)\n"
        "   - Testing activities (unit, integration, UAT)\n"
        "   - Deployment and go-live activities\n"
        "   - Post-deployment support activities\n\n"
        "**Project Summary Guidelines:**\n"
        "- `executive_summary`: 2-3 paragraph high-level summary of project goals, scope, and expected outcomes\n"
        "- `key_deliverables`: List 5-8 major deliverables (e.g., 'Fully functional mobile app', 'REST API with documentation')\n"
        "- `success_criteria`: List 4-6 measurable success metrics (e.g., 'System handles 10k concurrent users', 'API response time < 200ms')\n"
        '- `risks_and_mitigation`: List 4-6 project risks with mitigation strategies as objects with "risk" and "mitigation" fields (e.g., {"risk": "Third-party API downtime", "mitigation": "Implement fallback caching and retry logic"})\n\n'
        "**CRITICAL: Output ONLY the schema above. Do NOT add:**\n"
        "- ‚ùå \"cost_projection\" field (this will be auto-generated from resourcing_plan)\n"
        "- ‚ùå Any other fields not listed in the schema above\n"
        "- ‚ùå No markdown, no commentary, no explanations ‚Äî ONLY valid JSON matching the schema\n\n"
        "Scheduling Rules: \n"
        f"- The first activity must always start today ({today_str}).\n"
        "- If two activities are **independent**, overlap their timelines by **70‚Äì80%** of their duration (not full overlap)."
        "- If one activity **depends** on another, allow a small overlap of **10-15%** near the end of the predecessor if feasible."
        "- Avoid full serialization unless strictly required by dependency."
        "- Avoid full parallelism where all tasks start together ‚Äî stagger independent ones by **5-10%**."
        "- Ensure overall project duration stays **‚â§ 12 months**."
        "- Auto-calculate **End Date = Start Date + Effort Months**.\n"
        "- Auto-calculate **overview.Duration** as the total span in months from the earliest Start Date to the latest End Date.\n"
        "- `Complexity` should be simple, medium, or high based on duration of project.\n"
        "- **Always assign at least one Resource**."
        "- Distinguish `Owner` (responsible lead role) and `Resources` (supporting roles)."
        "\n"
        "**Critical: Owner and Resources Assignment Rules:**\n"
        "- `Owner` must ALWAYS be a valid JOB ROLE from the company's rate card.\n"
        "- `Owner` is NEVER an activity name, activity description, or task name.\n"
        "- `Resources` must contain only valid JOB ROLES from the company's rate card.\n"
        "- You MUST use ONLY the roles listed below - DO NOT invent new roles.\n"
        "- If `Resources` is missing, fallback to the same `Owner` role.\n"
        "- Use less resources as much as possible.\n"
        "\n"
        f"**üî¥ MANDATORY: Use ONLY these exact roles from the company's rate card:**\n"
        f"{chr(10).join('  - ' + role for role in (rate_card_roles or []))}\n"
        "\n"
        "**Examples of CORRECT Owner assignment:**\n"
        f"  ‚úì Owner: \"{rate_card_roles[0] if rate_card_roles else 'Backend Developer'}\" (this is a role from the rate card)\n"
        f"  ‚úì Owner: \"{rate_card_roles[1] if len(rate_card_roles) > 1 else 'Data Engineer'}\" (this is a role from the rate card)\n"
        f"  ‚úì Owner: \"{rate_card_roles[2] if len(rate_card_roles) > 2 else 'Solution Architect'}\" (this is a role from the rate card)\n"
        "\n"
        "**Examples of INCORRECT Owner assignment (DO NOT DO THIS):**\n"
        "  ‚úó Owner: \"Infrastructure Setup\" (this is an activity, not a role!)\n"
        "  ‚úó Owner: \"Data Ingestion Development\" (this is an activity, not a role!)\n"
        "  ‚úó Owner: \"Source Analysis\" (this is an activity, not a role!)\n"
        "  ‚úó Owner: \"John Smith\" (this is a person's name, not a role!)\n"
        "  ‚úó Owner: \"Project Manager\" (this role is NOT in the company's rate card!)\n"
        "\n"
        "Activity Duration Guidelines:\n"
        "Estimate realistic durations based on activity type and complexity. Use these as reference:\n"
        "\n"
        "**Planning & Design Activities:**\n"
        "- Requirements Gathering & Analysis: 0.5-1 month\n"
        "- System Architecture Design: 0.5-1 month\n"
        "- UI/UX Design & Wireframing: 0.75-1.5 months\n"
        "- Database Schema Design: 0.25-0.5 month\n"
        "- API Design & Documentation: 0.25-0.5 month\n"
        "\n"
        "**Development Activities:**\n"
        "- Simple CRUD Operations: 0.5-0.75 month\n"
        "- Complex Feature Development: 1-1.5 months\n"
        "- API Development (REST/GraphQL): 0.75-1.25 months\n"
        "- Database Implementation: 0.5-1 month\n"
        "- Authentication & Authorization: 0.75-1.25 months\n"
        "- Payment Gateway Integration: 1-1.5 months\n"
        "- Third-Party API Integrations: 0.5-1 month\n"
        "- Real-time Features (WebSockets, etc.): 1-1.5 months\n"
        "- Search Functionality: 0.75-1.25 months\n"
        "- File Upload/Management: 0.5-0.75 month\n"
        "- Notification System: 0.75-1 month\n"
        "- Reporting & Analytics: 1-1.5 months\n"
        "\n"
        "**AI/ML & Advanced Features:**\n"
        "- AI Model Integration: 1.5-2 months\n"
        "- Machine Learning Pipeline: 1.5-2.5 months\n"
        "- Natural Language Processing: 1.5-2 months\n"
        "- Computer Vision Features: 1.5-2 months\n"
        "- Recommendation Engine: 1-1.5 months\n"
        "\n"
        "**Testing & Quality Assurance:**\n"
        "- Unit Testing: 0.25-0.5 month\n"
        "- Integration Testing: 0.5-0.75 month\n"
        "- End-to-End Testing: 0.5-1 month\n"
        "- Performance Testing: 0.5-0.75 month\n"
        "- Security Testing: 0.75-1 month\n"
        "- User Acceptance Testing: 0.5-0.75 month\n"
        "\n"
        "**DevOps & Deployment:**\n"
        "- CI/CD Pipeline Setup: 0.5-0.75 month\n"
        "- Cloud Infrastructure Setup: 0.75-1 month\n"
        "- Containerization (Docker/K8s): 0.5-1 month\n"
        "- Monitoring & Logging Setup: 0.5-0.75 month\n"
        "- Production Deployment: 0.25-0.5 month\n"
        "\n"
        "**Domain-Specific Activity Templates:**\n"
        "\n"
        "**E-Commerce Domain:**\n"
        "- Product Catalog Management: 1-1.5 months\n"
        "- Shopping Cart & Checkout: 1.25-1.75 months\n"
        "- Order Management System: 1-1.5 months\n"
        "- Inventory Management: 1-1.5 months\n"
        "- Payment Processing: 1-1.5 months\n"
        "- Shipping Integration: 0.75-1 month\n"
        "\n"
        "**Healthcare Domain:**\n"
        "- Patient Management System: 1.5-2 months\n"
        "- Electronic Health Records (EHR): 2-2.5 months\n"
        "- Appointment Scheduling: 1-1.5 months\n"
        "- Medical Billing: 1.5-2 months\n"
        "- HIPAA Compliance Implementation: 1-1.5 months\n"
        "- Telemedicine Features: 1.5-2 months\n"
        "\n"
        "**FinTech Domain:**\n"
        "- Account Management: 1.5-2 months\n"
        "- Transaction Processing: 1.5-2 months\n"
        "- KYC/AML Compliance: 1.5-2 months\n"
        "- Fraud Detection System: 1.5-2.5 months\n"
        "- Financial Reporting: 1-1.5 months\n"
        "- Multi-Currency Support: 1-1.5 months\n"
        "\n"
        "**Education Domain:**\n"
        "- Learning Management System (LMS): 2-2.5 months\n"
        "- Course Management: 1-1.5 months\n"
        "- Student Portal: 1-1.5 months\n"
        "- Assessment & Grading: 1-1.5 months\n"
        "- Video Streaming Integration: 1-1.5 months\n"
        "- Certificate Generation: 0.5-0.75 month\n"
        "\n"
        "**Social Media/Community Domain:**\n"
        "- User Profiles & Authentication: 1-1.5 months\n"
        "- Feed/Timeline System: 1.5-2 months\n"
        "- Content Posting & Sharing: 1-1.5 months\n"
        "- Messaging/Chat System: 1.5-2 months\n"
        "- Notifications System: 0.75-1.25 months\n"
        "- Content Moderation: 1-1.5 months\n"
        "\n"
        "**IoT/Smart Systems Domain:**\n"
        "- Device Management: 1.5-2 months\n"
        "- Real-time Data Processing: 1.5-2 months\n"
        "- Sensor Data Analytics: 1.5-2 months\n"
        "- Remote Control Interface: 1-1.5 months\n"
        "- Alert & Automation System: 1-1.5 months\n"
        "\n"
        "**General Guidelines:**\n"
        "- For simple projects: Use lower end of duration ranges\n"
        "- For medium projects: Use mid-range durations\n"
        "- For complex projects: Use upper end or slightly beyond ranges\n"
        "- Activities can be split into smaller sub-activities if duration exceeds 2 months\n"
        "- Total project duration should realistically reflect the sum of critical path activities\n"
        "- Consider dependencies when scheduling - dependent activities should account for handoff time\n"
        "\n"
        "**CRITICAL: Infrastructure & Setup Activities - Use SHORT Durations!**\n"
        "Infrastructure and environment setup tasks are typically QUICK (1-2 weeks, NOT 1 month):\n"
        "- Azure/AWS/Cloud Infrastructure Setup: 0.25-0.5 month (1-2 weeks)\n"
        "- Database Environment Setup: 0.25-0.5 month (1-2 weeks)\n"
        "- CI/CD Pipeline Configuration: 0.25-0.5 month (1-2 weeks)\n"
        "- Development Environment Setup: 0.25 month (1 week)\n"
        "- Kubernetes/Container Setup: 0.5 month (2 weeks)\n"
        "- Monitoring & Logging Tools Setup: 0.25-0.5 month (1-2 weeks)\n"
        "\n"
        "**IMPORTANT: Use Granular Durations - NOT Everything is 1 Month!**\n"
        "Use realistic, varied durations based on actual effort required:\n"
        "- 0.25 month = 1 week (quick setup, configuration, simple tasks)\n"
        "- 0.5 month = 2 weeks (moderate complexity, integration work)\n"
        "- 0.75 month = 3 weeks (moderate to complex features)\n"
        "- 1 month = 4 weeks (complex features, major development)\n"
        "- 1.25-1.5 months = 5-6 weeks (very complex features, multiple integrations)\n"
        "- 1.75-2 months = 7-8 weeks (large system components, AI/ML work)\n"
        "\n"
        "**Activity Duration Examples (Realistic Estimates):**\n"
        "\n"
        "Example 1 - Infrastructure Setup:\n"
        "  Activity: \"Set up Azure infrastructure with SQL DB and monitoring\"\n"
        "  Duration: 0.5 month (2 weeks) ‚úì\n"
        "  NOT: 1 month ‚úó\n"
        "\n"
        "Example 2 - Data Source Analysis:\n"
        "  Activity: \"Analyze 30+ data sources and define integration approach\"\n"
        "  Duration: 0.75 month (3 weeks) ‚úì\n"
        "  NOT: 1 month ‚úó\n"
        "\n"
        "Example 3 - Simple ETL Pipeline:\n"
        "  Activity: \"Develop ETL pipeline for SQL database ingestion\"\n"
        "  Duration: 0.75 month (3 weeks) ‚úì\n"
        "  NOT: 1 month ‚úó\n"
        "\n"
        "Example 4 - Complex Feature:\n"
        "  Activity: \"Implement ML-based fraud detection system\"\n"
        "  Duration: 1.5-2 months (6-8 weeks) ‚úì\n"
        "\n"
        "Example 5 - Testing Phase:\n"
        "  Activity: \"Execute end-to-end testing and UAT\"\n"
        "  Duration: 0.5 month (2 weeks) ‚úì\n"
        "  NOT: 1 month ‚úó\n"
        "\n"
        "**Remember:** Most activities take LESS than 1 month! Use 0.25, 0.5, 0.75 frequently!\n"
        "\n"
        "- IDs must start from 1 and increment sequentially.\n"
        "- If the RFP or Knowledge Base text lacks detail, infer the missing pieces logically."
        "- Include all relevant roles and activities that ensure delivery of the project scope."
        "- Keep all field names exactly as in the schema.\n"
        f"{user_context}"
        f"RFP / Project Files Content:\n{rfp_text}\n\n"
        f"Knowledge Base Context (for enrichment only):\n{kb_context}\n"
        f"Clarification Q&A (User-confirmed answers take highest priority)\n"
        f"Use these answers to override or clarify any ambiguous or conflicting information.\n"
        f"Do NOT hallucinate beyond these facts.\n\n"
        f"{questions_context}\n\n"
        "‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ\n"
        "üö® FINAL CRITICAL REQUIREMENTS - READ THIS CAREFULLY:\n"
        "‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ‚îÅ\n"
        "1. ‚úÖ MUST include 'overview' object with all 7 fields\n"
        "2. ‚úÖ MUST include 'activities' array with at least 8-15 activities (THIS IS MANDATORY!)\n"
        "3. ‚úÖ MUST include 'resourcing_plan' as empty array []\n"
        "4. ‚úÖ MUST include 'project_summary' object with all 4 fields\n"
        "5. ‚ùå DO NOT return only metadata without activities\n"
        "6. ‚ùå DO NOT nest activities inside 'phases' - put them directly in 'activities' array\n"
        "7. ‚ùå DO NOT wrap the response in 'data' or 'project' keys - use the exact schema above\n"
        "8. üéØ Your response MUST be valid JSON that starts with '{' and ends with '}'\n\n"
        "REMEMBER: Output ONLY the JSON object. No explanations, no thinking, no markdown, no prose. Start your response with '{' and end with '}'. Nothing else.\n"
    )


def _build_questionnaire_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    """
    Build a prompt that forces the model to infer categories dynamically from RFP context.
    """
    name = getattr(project, "name", "Unnamed Project")
    domain = getattr(project, "domain", "General")
    tech = getattr(project, "tech_stack", "Modern Web Stack")
    compliance = getattr(project, "compliance", "General")
    duration = getattr(project, "duration", "TBD")

    return f"""
You are a **senior business analyst** preparing a requirement-clarification questionnaire
based on an RFP document.

Your goal: identify the main THEMES and subareas discussed in the RFP or Knowledge Base,
and then create **categories of questions** that align with those themes.
Do NOT reuse example categories blindly ‚Äî derive them from the content itself.

---

### Project Context
- Project Name: {name}
- Domain: {domain}
- Tech Stack: {tech}
- Compliance: {compliance}
- Duration: {duration}

### RFP Content
{rfp_text}

### Knowledge Base Context
{kb_chunks}

---

### TASK
1. First, analyze the RFP text to identify **key themes or topics** (e.g., Data Governance, SOX Controls,
   Cloud Migration, AI Enablement, Supply Chain Optimization, etc.).
2. For each theme, create a **category** with 5-6 specific questions.
3. Questions should clarify requirements, assumptions, or current-state processes.
4. Avoid repeating generic categories like "Architecture" or "Data & Security"
   unless they are explicitly discussed in the RFP.

---

### OUTPUT FORMAT
Return ONLY valid JSON in this structure:

{{
  "questions": [
    {{
      "category": "Data Governance & Ownership",
      "items": [
        {{
          "question": "Is there a defined data ownership model for finance data?",
          "user_understanding": "",
          "comment": ""
        }},
        {{
          "question": "Do you maintain audit logs for data corrections?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }},
    {{
      "category": "Regulatory Readiness and SOX Scope",
      "items": [
        {{
          "question": "What parts of the organization are in SOX scope?",
          "user_understanding": "",
          "comment": ""
        }}
      ]
    }}
  ]
}}

### RULES
- Categories must emerge logically from the RFP and KB text.
- Each category must contain at least 2 context-relevant questions.
- Each question must be concise, unambiguous, and require a short descriptive answer.
- Always include empty strings for 'user_understanding' and 'comment'.
- Output ONLY valid JSON (no explanations or markdown).
"""

def _extract_questions_from_text(raw_text: str) -> list[dict]:
    try:
        parsed = _extract_json(raw_text)

        # Case 1: Proper JSON with nested categories
        if isinstance(parsed, dict) and "questions" in parsed:
            qdata = parsed["questions"]
            if isinstance(qdata, list) and all(isinstance(x, dict) for x in qdata):
                # check if already nested structure
                if "items" in qdata[0]:
                    normalized = []
                    for cat in qdata:
                        normalized.append({
                            "category": cat.get("category", "General"),
                            "items": [
                                {
                                    "question": i.get("question", ""),
                                    "user_understanding": i.get("user_understanding", ""),
                                    "comment": i.get("comment", "")
                                } for i in cat.get("items", [])
                            ]
                        })
                    return normalized

                # Otherwise, flat ‚Üí group by category
                grouped = {}
                for q in qdata:
                    cat = q.get("category", "General") if isinstance(q, dict) else "General"
                    que = q.get("question", q) if isinstance(q, dict) else str(q)
                    grouped.setdefault(cat, []).append({
                        "question": que,
                        "user_understanding": "",
                        "comment": ""
                    })
                return [{"category": c, "items": lst} for c, lst in grouped.items()]

        # Case 2: List of plain questions
        if isinstance(parsed, list):
            return [{
                "category": "General",
                "items": [{"question": str(q), "user_understanding": "", "comment": ""} for q in parsed]
            }]
    except Exception:
        pass

    # Fallback ‚Äî parse raw text
    current_cat = "General"
    grouped: dict[str, list] = {}
    for line in raw_text.splitlines():
        line = line.strip()
        if not line:
            continue
        if re.match(r"^(#+\s*)?([A-Z][A-Za-z\s&/]+):?$", line) and not line.endswith("?"):
            current_cat = re.sub(r"^#+\s*", "", line).strip(": ").strip()
            continue
        if "?" in line:
            qtext = re.sub(r"^\d+[\).\s]+", "", line).strip()
            grouped.setdefault(current_cat, []).append({
                "question": qtext,
                "user_understanding": "",
                "comment": ""
            })

    return [{"category": c, "items": lst} for c, lst in grouped.items()]
async def generate_project_questions(db: AsyncSession, project) -> dict:
    """
    Generate a categorized questionnaire for the given project using Ollama.
    Saves the questions.json file in Azure Blob.
    """

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        if getattr(project, "files", None):
            files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            if files:
                rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"Failed to extract RFP for questions: {e}")

    # ---------- Retrieve Knowledge Base ----------
    kb_results = _rag_retrieve(rfp_text or project.name or project.domain)
    kb_chunks = [ch["content"] for group in kb_results for ch in group["chunks"]] if kb_results else []

    # ---------- Build prompt ----------
    prompt = _build_questionnaire_prompt(rfp_text, kb_chunks, project)

    # ---------- Query Ollama ----------
    try:
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.8))
        questions = _extract_questions_from_text(raw_text)
        total_q = sum(len(cat["items"]) for cat in questions)
        logger.info(f" Generated {total_q} questions under {len(questions)} categories for project {project.id}")

        # ---------- Save to Blob Storage ----------
        blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        try:
            await azure_blob.upload_bytes(
                json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
            )

            db_file = models.ProjectFile(
                project_id=project.id,
                file_name="questions.json",
                file_path=blob_name,
            )

            db.add(db_file)
            await db.commit()
            await db.refresh(db_file)

            logger.info(f" Saved questions.json for project {project.id}")
        except Exception as e:
            logger.warning(f"Failed to upload questions.json: {e}")

        return {"questions": questions}

    except Exception as e:
        logger.error(f" Question generation failed: {e}")
        return {"questions": []}
    
# Update questions.json with user input answers
async def update_questions_with_user_input(
    db: AsyncSession, project, user_answers: dict
) -> dict:
    from app.utils import azure_blob

    blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
    try:
        # Load current questions.json
        q_bytes = await azure_blob.download_bytes(blob_name)
        q_json = json.loads(q_bytes.decode("utf-8"))
        questions = q_json.get("questions", [])

        # Merge answers into the structure
        for cat in questions:
            cat_name = cat.get("category")
            for item in cat.get("items", []):
                q_text = item.get("question")
                ans = (
                    user_answers.get(cat_name, {}).get(q_text)
                    if user_answers.get(cat_name)
                    else None
                )
                if ans:
                    item["user_understanding"] = ans

        # Upload updated JSON to Blob
        new_bytes = json.dumps({"questions": questions}, ensure_ascii=False, indent=2).encode("utf-8")
        await azure_blob.upload_bytes(new_bytes, blob_name)
        logger.info(f" Updated questions.json with user input for project {project.id}")

        #  Save / update DB record
        db_file = models.ProjectFile(
            project_id=project.id,
            file_name="questions.json",
            file_path=blob_name,
        )
        db.add(db_file)
        await db.commit()
        await db.refresh(db_file)

        return {"questions": questions}

    except Exception as e:
        logger.error(f"Failed to update questions.json with user input: {e}")
        return {}

    
def _build_architecture_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    name = (getattr(project, "name", "") or "Untitled Project").strip()
    domain = (getattr(project, "domain", "") or "General").strip()
    tech = (getattr(project, "tech_stack", "") or "Modern Web + Cloud Stack").strip()

    return f"""
    You are a **senior enterprise solution architect** tasked with designing a *tailored cloud system architecture diagram*
    strictly based on the provided RFP and contextual knowledge.

    ### PROJECT CONTEXT
    - **Project Name:** {name}
    - **Domain:** {domain}
    - **Tech Stack:** {tech}

    ### RFP SUMMARY
    {rfp_text}

    ### KNOWLEDGE BASE CONTEXT
    {kb_chunks}

    ---

    ###  STEP 1 ‚Äî Reasoning (Internal)
    Analyze the provided RFP and knowledge base to:
    1. Identify all domain-specific **entities, systems, or technologies** mentioned or implied.
    2. Categorize each component into the most appropriate architecture layer:
    - Frontend (UI/Apps)
    - Backend (Services/APIs)
    - Data (Databases, Storage, External APIs)
    - AI/Analytics (ML, Insights, NLP, Recommendations)
    - Security/Monitoring/DevOps (IAM, Key Vault, CI/CD, Logging)
    3. Infer **connections and data flows** between components (e.g., API requests, pipelines, message queues).
    4. Skip any layers not relevant to this RFP.

    You will use this reasoning to build the architecture ‚Äî but **do not include this reasoning** in your final output.

    ---

    ###  STEP 2 ‚Äî Graphviz DOT Output
    Generate **only valid Graphviz DOT code** representing the inferred architecture.

    Follow these rules strictly:
    - Begin with: `digraph Architecture {{`
    - End with: `}}`
    - Use **horizontal layout** ‚Üí `rankdir=LR`
    - Include **only relevant clusters** (omit unused layers)
    - Keep ‚â§ 15 nodes total
    - Use **orthogonal edges** (`splines=ortho`)
    - Each node label must clearly represent an actual system, service, or tool
    - Logical flow should follow Frontend ‚Üí Backend ‚Üí Data ‚Üí AI ‚Üí Security (only if applicable)
    -  **Ensure data layers both receive and provide information** ‚Äî show arrows *into* and *out of* data/storage nodes if analytics, AI, or reporting components exist.

    ---

    ### VISUAL STYLE
    - **Graph:** dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.3
    - **Clusters:** style="filled,rounded", fontname="Helvetica-Bold", fontsize=13
    - **Node Shapes and Colors:**
    - Frontend ‚Üí `box`, pastel blue (`fillcolor="#E3F2FD"`)
    - Backend/API ‚Üí `box3d`, pastel green (`fillcolor="#E8F5E9"`)
    - Data/Storage ‚Üí `cylinder`, pastel yellow (`fillcolor="#FFFDE7"`)
    - AI/Analytics ‚Üí `ellipse`, pastel purple (`fillcolor="#F3E5F5"`)
    - Security/Monitoring ‚Üí `diamond`, gray (`fillcolor="#ECEFF1"`)
    - **Edges:** color="#607D8B", penwidth=1.5, arrowsize=0.9

    ---

    ###  STEP 3 ‚Äî Domain Intelligence (Auto-Enrichment)
    If applicable, automatically enrich the architecture using these domain patterns:

    - **FinTech** ‚Üí Payment Gateway, Fraud Detection, KYC/AML Service, Ledger DB
    - **HealthTech** ‚Üí Patient Portal, EHR System, FHIR API, HIPAA Compliance Layer
    - **GovTech** ‚Üí Citizen Portal, Secure API Gateway, Compliance & Audit Logging
    - **AI/ML Projects** ‚Üí Model API, Embedding Store, Training Pipeline, Monitoring Service
    - **Data Platforms** ‚Üí ETL Pipeline, Data Lake, BI Dashboard
    - **Enterprise SaaS** ‚Üí Tenant Manager, Auth Service, Billing & Subscription Module

    Include these elements **only if they logically fit** the RFP description.

    ---

    ###  STEP 4 ‚Äî OUTPUT RULES (CRITICAL!)

    **YOUR RESPONSE MUST START WITH:** digraph Architecture {{
    **YOUR RESPONSE MUST END WITH:** }}

    - Output *only* valid Graphviz DOT syntax
    - **NO** markdown code fences
    - **NO** explanatory text before or after the DOT code
    - **NO** reasoning or commentary
    - **NO** sentences like "Based on the analysis..." or "Here is the code..."
    - **NO** C-style comments (//) - DOT does not support them! Use # or /* */ if needed
    - **NO** escaped quotes (\") - Use plain quotes in attribute values
    - **IMPORTANT:** Ensure all labels are complete and properly closed with quotes
    - **IMPORTANT:** Do NOT end labels with colons (e.g., label="Database:" is WRONG, use label="Database")
    - The FIRST character of your response must be "d" (from digraph)
    - The LAST character of your response must be closing brace

    **WRONG (Do NOT do this):**
    Based on the analysis, here is the code:
    digraph Architecture {{ ... }}

    **ALSO WRONG - Incomplete labels:**
    SQL_DB [label="SQL Database\nTransformation:

    **CORRECT (Do this):**
    digraph Architecture {{ ... }}
    SQL_DB [label="SQL Database\nTransformation"]

    Your response must be pure DOT code that can be directly passed to Graphviz without any processing.
    """


def _build_eraser_architecture_prompt(rfp_text: str, kb_chunks: List[str], project=None) -> str:
    """
    Build prompt for Eraser.io DSL architecture diagram generation.
    Focuses on using EXACT tech stack from RFP.
    """
    name = (getattr(project, "name", "") or "Untitled Project").strip()
    domain = (getattr(project, "domain", "") or "General").strip()
    tech = (getattr(project, "tech_stack", "") or "Modern Web + Cloud Stack").strip()

    # Convert tech_stack string to list if needed
    tech_list = []
    if tech:
        if isinstance(tech, str):
            # Split by common delimiters
            tech_list = [t.strip() for t in re.split(r'[,;|]', tech) if t.strip()]
        elif isinstance(tech, list):
            tech_list = tech

    tech_list_str = "\n".join(f"  - {t}" for t in tech_list) if tech_list else "  - (No specific tech stack provided)"

    return f"""
    You are a **senior cloud architect** creating an **Eraser.io architecture diagram**.

    ### PROJECT CONTEXT
    - **Project Name:** {name}
    - **Domain:** {domain}
    - **Tech Stack (CRITICAL - USE THESE EXACT TECHNOLOGIES):**
{tech_list_str}

    ### RFP SUMMARY
    {rfp_text}

    ### KNOWLEDGE BASE CONTEXT
    {kb_chunks}

    ---

    ### TASK
    Generate **Eraser.io DSL syntax** for a cloud architecture diagram.

    **CRITICAL REQUIREMENTS:**
    1. **USE ONLY THE TECH STACK LISTED ABOVE** - Do NOT invent or add technologies not in the tech stack
    2. Each technology from the tech stack MUST appear as a node in the diagram
    3. Use appropriate cloud icons for each technology
    4. Show logical data flows and connections
    5. Group related components together

    ---

    ### ERASER.IO DSL SYNTAX RULES

    **Nodes:**
    ```
    NodeName [icon: icon-name, color: color-name]
    ```

    **Groups (containers):**
    ```
    GroupName {{
      Node1 [icon: aws-lambda]
      Node2 [icon: aws-s3]
    }}
    ```

    **Connections (arrows):**
    ```
    Node1 > Node2
    Node1 > Node2, Node3, Node4
    ```

    **Available Cloud Icons:**
    - **Azure:** azure-functions, azure-blob-storage, azure-sql-database, azure-cosmos-db, azure-app-service, azure-api-management, azure-data-factory, azure-databricks, azure-synapse-analytics, azure-power-bi, azure-devops, azure-kubernetes-service, azure-virtual-machines
    - **AWS:** aws-lambda, aws-s3, aws-rds, aws-dynamodb, aws-ec2, aws-api-gateway, aws-ecs, aws-eks, aws-cloudfront, aws-sqs, aws-sns
    - **GCP:** gcp-cloud-functions, gcp-cloud-storage, gcp-cloud-sql, gcp-firestore, gcp-compute-engine, gcp-kubernetes-engine
    - **General:** database, server, cloud, api, monitor, tool, globe

    ---

    ### DOMAIN-SPECIFIC PATTERNS (Use if matching domain)

    - **Data Analytics/BI:** ETL Pipeline, Data Lake, Data Warehouse, BI Dashboard, Analytics Engine
    - **FinTech:** Payment Gateway, Fraud Detection, KYC Service, Transaction DB, Ledger
    - **HealthTech:** Patient Portal, EHR System, FHIR API, Compliance Layer
    - **AI/ML:** Model API, Training Pipeline, Feature Store, Model Registry
    - **E-Commerce:** Product Catalog, Shopping Cart, Payment Processor, Order Management

    ---

    ### OUTPUT RULES (CRITICAL!)

    **YOUR RESPONSE MUST:**
    1. Start immediately with node/group definitions (no explanatory text)
    2. Use ONLY technologies from the tech stack provided above
    3. Be pure Eraser.io DSL syntax
    4. NOT include markdown, commentary, or explanations
    5. Map each tech stack item to appropriate cloud icon

    **WRONG (Do NOT do this):**
    ```
    Based on the analysis, here's the architecture:
    VPC {{ ... }}
    ```

    **CORRECT (Do this):**
    ```
    Cloud Infrastructure {{
      Azure Data Factory [icon: azure-data-factory, color: blue]
      Azure Databricks [icon: azure-databricks, color: orange]
    }}

    Azure Data Factory > Azure Databricks
    ```

    **TECH STACK MAPPING EXAMPLES:**
    - "Azure Data Factory" ‚Üí `Azure Data Factory [icon: azure-data-factory]`
    - "Power BI" ‚Üí `Power BI Dashboard [icon: azure-power-bi]`
    - "Azure SQL Database" ‚Üí `Azure SQL DB [icon: azure-sql-database]`
    - "Kubernetes" ‚Üí `Kubernetes Cluster [icon: azure-kubernetes-service]`
    - "React" ‚Üí `React Frontend [icon: react]`
    - "Node.js" ‚Üí `Node.js API [icon: nodejs]`

    **Remember:** Your output must be **PURE Eraser.io DSL** with NO additional text!
    """


async def _call_eraser_api(dsl_code: str) -> tuple[str | None, str | None]:
    """
    Call Eraser.io API to render architecture diagram.
    Returns: (image_url, editor_url) tuple or (None, None) on failure
    """
    from app.config.config import ERASER_IO_API_KEY, ERASER_IO_API_URL

    if not ERASER_IO_API_KEY:
        logger.warning("‚ö†Ô∏è ERASER_IO_API_KEY not configured - skipping Eraser.io diagram generation")
        return None, None

    headers = {
        "Authorization": f"Bearer {ERASER_IO_API_KEY}",
        "Content-Type": "application/json"
    }

    payload = {
        "theme": "light",
        "background": True,
        "elements": [
            {
                "type": "diagram",
                "diagramType": "cloud-architecture-diagram",
                "code": dsl_code
            }
        ]
    }

    try:
        logger.info(f"üé® Calling Eraser.io API to render architecture diagram...")
        response = await anyio.to_thread.run_sync(
            lambda: requests.post(ERASER_IO_API_URL, headers=headers, json=payload, timeout=30)
        )

        if response.status_code == 200:
            result = response.json()
            image_url = result.get("imageUrl")
            editor_url = result.get("createEraserFileUrl")
            logger.info(f"‚úÖ Eraser.io diagram generated successfully: {image_url}")
            return image_url, editor_url
        else:
            logger.error(f"‚ùå Eraser.io API error: {response.status_code} - {response.text}")
            return None, None

    except Exception as e:
        logger.error(f"‚ùå Eraser.io API call failed: {e}")
        return None, None


async def generate_architecture_eraser(
    db: AsyncSession,
    project,
    rfp_text: str,
    kb_chunks: List[str],
    blob_base_path: str,
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate architecture diagram using Eraser.io API.
    Downloads PNG from Eraser.io and stores in Azure Blob.
    Falls back to Graphviz if Eraser.io is not configured or fails.
    """
    from app.config.config import ERASER_IO_API_KEY

    # Check if Eraser.io is configured
    if not ERASER_IO_API_KEY:
        logger.info("üìä Eraser.io not configured - using Graphviz fallback")
        return await generate_architecture(db, project, rfp_text, kb_chunks, blob_base_path)

    prompt = _build_eraser_architecture_prompt(rfp_text, kb_chunks, project)

    # Step 1: Generate Eraser.io DSL from LLM
    async def _generate_dsl_from_ai(retry: int = 0) -> str:
        """Call Ollama to generate Eraser.io DSL."""
        try:
            return await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.7))
        except Exception as e:
            if retry < 2:
                logger.warning(f"Ollama call failed (retry {retry+1}/3): {e}")
                await anyio.sleep(2)
                return await _generate_dsl_from_ai(retry + 1)
            logger.error(f"Ollama DSL generation failed after retries: {e}")
            return ""

    dsl_code = await _generate_dsl_from_ai()
    if not dsl_code:
        logger.warning("‚ö†Ô∏è No DSL code returned by AI - using Graphviz fallback")
        return await generate_architecture(db, project, rfp_text, kb_chunks, blob_base_path)

    # Step 2: Clean DSL code
    dsl_code = re.sub(r"```[a-zA-Z]*", "", dsl_code).replace("```", "").strip()
    dsl_code = dsl_code.strip("`").strip()

    logger.info(f"üìù Generated Eraser.io DSL ({len(dsl_code)} chars)")

    # Step 3: Call Eraser.io API
    image_url, editor_url = await _call_eraser_api(dsl_code)

    if not image_url:
        logger.warning("‚ö†Ô∏è Eraser.io rendering failed - using Graphviz fallback")
        return await generate_architecture(db, project, rfp_text, kb_chunks, blob_base_path)

    # Step 4: Download PNG from Eraser.io
    try:
        logger.info(f"üì• Downloading diagram from Eraser.io: {image_url}")
        png_response = await anyio.to_thread.run_sync(
            lambda: requests.get(image_url, timeout=30)
        )
        png_response.raise_for_status()
        png_bytes = png_response.content

        # Step 5: Upload to Azure Blob
        blob_name_png = f"{blob_base_path}/architecture_eraser_{project.id}.png"
        await azure_blob.upload_bytes(blob_name_png, png_bytes, content_type="image/png")
        logger.info(f"‚úÖ Uploaded Eraser.io diagram to Azure: {blob_name_png}")

        # Step 6: Store in database
        db_file = models.ProjectFile(
            project_id=project.id,
            file_name=f"architecture_eraser_{project.id}.png",
            file_path=blob_name_png,
            file_type="image/png",
        )
        db.add(db_file)
        await db.commit()
        await db.refresh(db_file)

        logger.info(f"‚úÖ Eraser.io architecture diagram stored for project {project.id}: {blob_name_png}")
        return db_file, blob_name_png

    except Exception as e:
        logger.error(f"‚ùå Failed to download/store Eraser.io diagram: {e}")
        logger.info("‚ö†Ô∏è Falling back to Graphviz")
        return await generate_architecture(db, project, rfp_text, kb_chunks, blob_base_path)


async def _generate_fallback_architecture(
    db: AsyncSession,
    project,
    blob_base_path: str
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate and upload a default fallback architecture diagram (4-layer generic layout).
    Triggered when Ollama or Graphviz generation fails.
    """
    logger.warning(" Using fallback default architecture layout")

    # --- Default DOT diagram ---
    fallback_dot = """
digraph Architecture {
    rankdir=LR;
    graph [dpi=200, bgcolor="white", nodesep=1.3, ranksep=1.2, splines=ortho];
    node [style="rounded,filled", fontname="Helvetica-Bold", fontsize=13, penwidth=1.2];

    subgraph cluster_frontend {
        label="Frontend / User Touchpoints";
        style="filled,rounded"; fillcolor="#E3F2FD";
        web[label="Web App (React / Angular)", shape=box, fillcolor="#BBDEFB"];
        mobile[label="Mobile App", shape=box, fillcolor="#BBDEFB"];
    }

    subgraph cluster_backend {
        label="Backend / Services";
        style="filled,rounded"; fillcolor="#E8F5E9";
        api[label="Core API (FastAPI / Node.js)", shape=box3d, fillcolor="#C8E6C9"];
        auth[label="Auth Service", shape=box3d, fillcolor="#C8E6C9"];
    }

    subgraph cluster_data {
        label="Data / Storage";
        style="filled,rounded"; fillcolor="#FFFDE7";
        db[label="Database (PostgreSQL)", shape=cylinder, fillcolor="#FFF9C4"];
        blob[label="Blob Storage", shape=cylinder, fillcolor="#FFF9C4"];
    }

    subgraph cluster_ai {
        label="AI / Analytics";
        style="filled,rounded"; fillcolor="#F3E5F5";
        ai[label="AI Engine / Insights", shape=ellipse, fillcolor="#E1BEE7"];
        dashboard[label="BI Dashboard", shape=ellipse, fillcolor="#E1BEE7"];
    }

    # Data flow (using xlabels to avoid orthogonal label warnings)
    web -> api [xlabel="HTTP Request"];
    mobile -> api [xlabel="Mobile API Call"];
    api -> db [xlabel="DB Query"];
    db -> ai [xlabel="ETL/Inference"];
    ai -> dashboard [xlabel="Visualization"];
    api -> auth [xlabel="Auth Validation"];

}
"""

    # --- Render DOT ‚Üí PNG & SVG ---
    tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
    try:
        graph = graphviz.Source(fallback_dot, engine="dot")
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Fallback Graphviz rendering failed: {e}")
        return None, ""

    # --- Upload both files to Azure Blob ---
    blob_name_png = f"{blob_base_path}/architecture_fallback_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_fallback_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)
        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # --- Save both records in DB ---
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Fallback architecture diagrams stored for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png



async def generate_architecture(
    db: AsyncSession,
    project,
    rfp_text: str,
    kb_chunks: List[str],
    blob_base_path: str,
) -> tuple[models.ProjectFile | None, str]:
    """
    Generate a visually clean, context-aware architecture diagram (PNG & SVG)
    from RFP + KB context using Ollama + Graphviz.
    Uses dynamic prompts that adapt layers automatically (no static template).
    Includes retry logic, sanitization, validation, and fallback diagram.
    """

    prompt = _build_architecture_prompt(rfp_text, kb_chunks, project)

    # ---------- Step 1: Ask Ollama for Graphviz DOT code ----------
    async def _generate_dot_from_ai(retry: int = 0) -> str:
        """Call Ollama locally to generate DOT diagram."""
        try:
            return await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.7))
        except Exception as e:
            if retry < 2:
                logger.warning(f"Ollama call failed (retry {retry+1}/3): {e}")
                await anyio.sleep(2)
                return await _generate_dot_from_ai(retry + 1)
            logger.error(f"Ollama architecture generation failed after retries: {e}")
            return ""


    dot_code = await _generate_dot_from_ai()
    if not dot_code:
        logger.warning(" No DOT code returned by AI ‚Äî generating fallback diagram")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 2: Clean & sanitize DOT ----------
    # Remove markdown code fences
    dot_code = re.sub(r"```[a-zA-Z]*", "", dot_code).replace("```", "").strip()
    dot_code = dot_code.strip("`").strip()

    # Extract only the DOT code if LLM added commentary
    # Look for "digraph" and extract from there to the last closing brace
    match = re.search(r'(digraph\s+\w+\s*\{.*\})\s*$', dot_code, re.DOTALL | re.IGNORECASE)
    if match:
        dot_code = match.group(1).strip()
    else:
        # Try to find any digraph block
        match = re.search(r'digraph\s+\w+\s*\{', dot_code, re.IGNORECASE)
        if match:
            # Extract from digraph to the end
            start_idx = match.start()
            dot_code = dot_code[start_idx:].strip()

    # Remove duplicate digraph declarations (common LLM error)
    # If we have multiple "digraph Architecture {" lines, keep only the first one
    lines = dot_code.split('\n')
    digraph_count = 0
    cleaned_lines = []
    for line in lines:
        if re.match(r'^\s*digraph\s+\w+\s*\{\s*$', line, re.IGNORECASE):
            digraph_count += 1
            if digraph_count == 1:
                cleaned_lines.append(line)
            # else: skip duplicate digraph lines
        else:
            cleaned_lines.append(line)
    dot_code = '\n'.join(cleaned_lines)

    dot_code = re.sub(r"(?i)^graph\s", "digraph ", dot_code)

    # Remove C-style comments (// ...) - Graphviz DOT doesn't support them
    dot_code = re.sub(r'//[^\n]*', '', dot_code)

    # Fix escaped quotes - DOT doesn't need escaped quotes in attribute values
    dot_code = dot_code.replace('\\"', '"')

    # Remove extra whitespace and blank lines
    dot_code = '\n'.join(line for line in dot_code.split('\n') if line.strip())

    # Fix brace mismatch
    open_braces = dot_code.count("{")
    close_braces = dot_code.count("}")
    if open_braces > close_braces:
        dot_code += "}" * (open_braces - close_braces)
    elif close_braces > open_braces:
        dot_code = "digraph Architecture {\n" + dot_code

    if not dot_code.lower().startswith("digraph"):
        dot_code = f"digraph Architecture {{\n{dot_code}\n}}"

    # Remove control characters
    dot_code = re.sub(r"[^\x09\x0A\x0D\x20-\x7E]", "", dot_code)

    # Fix label syntax - ensure all labels are properly quoted and escaped
    # Replace any labels that might have unescaped colons or special characters
    def fix_label(match):
        """Fix label syntax by properly escaping special characters."""
        label_content = match.group(1)
        # Remove any problematic colons at the end of labels (incomplete labels)
        label_content = re.sub(r':\s*$', '', label_content)
        # Escape quotes inside labels
        label_content = label_content.replace('"', '\\"')
        return f'label="{label_content}"'

    # Fix incomplete or malformed labels
    dot_code = re.sub(r'label="([^"]*)"?\s*(?=[;\]\}]|\n|$)', fix_label, dot_code)

    # Ensure node names don't have special characters that need escaping
    # Replace problematic node names with safe versions
    dot_code = re.sub(r'([A-Za-z0-9_]+)\s*\[label="([^"]+):"', r'\1 [label="\2"', dot_code)

    # ---------- Step 3: Do NOT override GPT's style ----------
    # Keep GPT's own clusters, nodes, and colors ‚Äî just ensure it's syntactically valid
    # (Old static preamble removed intentionally)

    # ---------- Step 4: Render DOT ‚Üí PNG & SVG ----------
    try:
        tmp_base = tempfile.NamedTemporaryFile(delete=False, suffix=".dot").name
        graph = graphviz.Source(dot_code, engine="dot")

        # Render both PNG and SVG for better clarity
        graph.render(tmp_base, format="png", cleanup=True)
        graph.render(tmp_base, format="svg", cleanup=True)

        png_path = tmp_base + ".png"
        svg_path = tmp_base + ".svg"
    except Exception as e:
        logger.error(f" Graphviz rendering failed: {e}\n--- DOT Snippet ---\n{dot_code[:800]}")
        return await _generate_fallback_architecture(db, project, blob_base_path)

    # ---------- Step 5: Upload PNG to Azure Blob ----------
    blob_name_png = f"{blob_base_path}/architecture_{project.id}.png"
    blob_name_svg = f"{blob_base_path}/architecture_{project.id}.svg"

    try:
        with open(png_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_png)

        with open(svg_path, "rb") as fh:
            await azure_blob.upload_bytes(fh.read(), blob_name_svg)
    finally:
        for path in [png_path, svg_path, tmp_base]:
            try:
                os.remove(path)
            except FileNotFoundError:
                pass

    # ---------- Step 6: Replace old architecture file ----------
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "architecture.png",
        )
    )
    old_file = result.scalars().first()
    if old_file:
        try:
            await azure_blob.delete_blob(old_file.file_path)
            await db.delete(old_file)
            await db.commit()
        except Exception as e:
            logger.warning(f" Failed to delete old architecture.png: {e}")

    # ---------- Step 7: Save new ProjectFile records (PNG + SVG) ----------
    db_file_png = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.png",
        file_path=blob_name_png,
    )
    db_file_svg = models.ProjectFile(
        project_id=project.id,
        file_name="architecture.svg",
        file_path=blob_name_svg,
    )

    db.add_all([db_file_png, db_file_svg])
    await db.commit()
    await db.refresh(db_file_png)
    await db.refresh(db_file_svg)

    logger.info(
        f" Architecture diagrams stored successfully for project {project.id}: "
        f"{blob_name_png}, {blob_name_svg}"
    )

    return db_file_png, blob_name_png


# --- Cleaner ---
async def clean_scope(db: AsyncSession, data: Dict[str, Any], project=None) -> Dict[str, Any]:
    if not isinstance(data, dict):
        return {}

    ist = pytz.timezone("Asia/Kolkata")
    # Use timezone-naive datetime to avoid comparison issues with parsed dates
    today = datetime.now(ist).replace(hour=0, minute=0, second=0, microsecond=0, tzinfo=None)

    activities: List[Dict[str, Any]] = []
    start_dates, end_dates = [], []
    role_month_map: Dict[str, Dict[str, float]] = {}
    role_order: List[str] = []

    # --- Helper: compute monthly allocation based on actual days in month ---
    def month_effort(s: datetime, e: datetime) -> Dict[str, float]:
        cur = s
        month_eff = {}
        while cur <= e:
            year, month = cur.year, cur.month
            days_in_month = monthrange(year, month)[1]
            start_day = cur.day if cur.month == s.month else 1
            end_day = e.day if cur.month == e.month else days_in_month
            days_count = end_day - start_day + 1
            month_eff[f"{cur.strftime('%b %Y')}"] = round(days_count / 30.0, 2)
            # move to next month
            if month == 12:
                cur = datetime(year + 1, 1, 1)
            else:
                cur = datetime(cur.year, cur.month + 1, 1)
        return month_eff

    # --- Process activities ---
    for idx, a in enumerate(data.get("activities") or [], start=1):
        owner = a.get("Owner") or "Unassigned"

        # Parse dependencies
        raw_deps = [d.strip() for d in str(a.get("Resources") or "").split(",") if d.strip()]

        # Remove owner from resources if duplicated
        raw_deps = [r for r in raw_deps if r.lower() != owner.lower()]

        # Owner always included, then other resources
        roles = [owner] + raw_deps

        s = _parse_date_safe(a.get("Start Date"), today)
        e = _parse_date_safe(a.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        # --- allocate per month (no splitting among roles) ---
        month_alloc = month_effort(s, e)
        for role in roles:
            if role not in role_month_map:
                role_month_map[role] = {}
                role_order.append(role)
            for m, eff in month_alloc.items():
                role_month_map[role][m] = role_month_map[role].get(m, 0.0) + eff

        dur_days = max(1, (e - s).days)
        activities.append({
            "ID": idx,
            "Activities": _safe_str(a.get("Activities")),
            "Description": _safe_str(a.get("Description")),
            "Owner": owner,
            "Resources": ", ".join(raw_deps),
            "Start Date": s,
            "End Date": e,
            "Effort Months": round_to_half(dur_days / 30.0),  # Round to nearest 0.5
        })

        start_dates.append(s)
        end_dates.append(e)

        # --- Sort activities ---
    activities.sort(key=lambda x: x["Start Date"])
    for idx, a in enumerate(activities, start=1):
        a["ID"] = idx
        a["Start Date"] = a["Start Date"].strftime("%Y-%m-%d")
        a["End Date"] = a["End Date"].strftime("%Y-%m-%d")

    # --- Project span & month labels (Month 1, Month 2, ...) ---
    min_start = min(start_dates) if start_dates else today
    max_end = max(end_dates) if end_dates else min_start
    duration = max(1.0, round(max(1, (max_end - min_start).days) / 30.0, 2))
    total_months = max(1, math.ceil((max_end - min_start).days / 30.0))

    month_labels = [f"Month {i}" for i in range(1, total_months + 1)]

    # --- Build per-role, per-month day usage ---
    role_month_usage: Dict[str, Dict[str, float]] = {r: {m: 0.0 for m in month_labels} for r in role_order}

    # Compute total active days per relative month window
    for act in activities:
        s = _parse_date_safe(act.get("Start Date"), today)
        e = _parse_date_safe(act.get("End Date"), s + timedelta(days=30))
        if e < s:
            e = s + timedelta(days=30)

        involved_roles = [act.get("Owner") or "Unassigned"] + [
            r.strip() for r in str(act.get("Resources") or "").split(",") if r.strip()
        ]

        for m_idx in range(total_months):
            rel_start = min_start + timedelta(days=m_idx * 30)
            rel_end = min_start + timedelta(days=(m_idx + 1) * 30)

            # overlap between activity and this relative month window
            overlap_start = max(s, rel_start)
            overlap_end = min(e, rel_end)
            overlap_days = 0
            if overlap_end >= overlap_start:
                overlap_days = (overlap_end - overlap_start).days + 1

            if overlap_days > 0:
                for r in involved_roles:
                    if r not in role_month_usage:
                        role_month_usage[r] = {ml: 0.0 for ml in month_labels}
                    role_month_usage[r][f"Month {m_idx + 1}"] += overlap_days

    # --- Convert days to effort with 4-tier partial-month logic ---
    for r, months in role_month_usage.items():
        for m, days in months.items():
            if days > 21:
                months[m] = 1.0
            elif 15 <= days <= 21:
                months[m] = 0.75
            elif 8 <= days < 15:
                months[m] = 0.5
            elif 1 <= days < 8:
                months[m] = 0.25
            else:
                months[m] = 0.0

    try:
        if db:
            ROLE_RATE_MAP_DYNAMIC = await get_rate_map_for_project(db, project)
        else:
            ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP
    except Exception as e:
        logger.warning(f"Rate map fallback due to error: {e}")
        ROLE_RATE_MAP_DYNAMIC = ROLE_RATE_MAP


    # --- Build final resourcing plan ---
    resourcing_plan = []
    for idx, role in enumerate(role_order, start=1):
        # Round each month's effort to nearest 0.5
        month_efforts_raw = role_month_usage.get(role, {m: 0 for m in month_labels})
        month_efforts = {m: round_to_half(float(v)) for m, v in month_efforts_raw.items()}

        # Calculate and round total effort
        total_effort = round_to_half(sum(month_efforts.values()))

        rate = ROLE_RATE_MAP_DYNAMIC.get(role, ROLE_RATE_MAP.get(role, 2000.0))
        cost = round(total_effort * rate, 2)
        plan_entry = {
            "ID": idx,
            "Resources": role,
            "Rate/month": rate,
            **month_efforts,
            "Efforts": total_effort,
            "Cost": cost,
        }
        resourcing_plan.append(plan_entry)

    # --- Apply discount if present ---
    discount_percentage = data.get("discount_percentage", 0)
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        discount_multiplier = 1 - (discount_percentage / 100.0)
        logger.info(f"üí∞ Applying {discount_percentage}% discount (multiplier: {discount_multiplier})")

        # Apply discount to all costs in resourcing_plan
        for plan_entry in resourcing_plan:
            original_cost = plan_entry.get("Cost", 0)
            discounted_cost = round(original_cost * discount_multiplier, 2)
            plan_entry["Cost"] = discounted_cost
            logger.info(f"  ‚Üí {plan_entry['Resources']}: ${original_cost} ‚Üí ${discounted_cost}")

    # --- Overview ---
    ov = data.get("overview") or {}
    data["overview"] = {
        "Project Name": _safe_str(ov.get("Project Name") or getattr(project, "name", "Untitled Project")),
        "Domain": _safe_str(ov.get("Domain") or getattr(project, "domain", "")),
        "Complexity": _safe_str(ov.get("Complexity") or getattr(project, "complexity", "")),
        "Tech Stack": _safe_str(ov.get("Tech Stack") or getattr(project, "tech_stack", "")),
        "Use Cases": _safe_str(ov.get("Use Cases") or getattr(project, "use_cases", "")),
        "Compliance": _safe_str(ov.get("Compliance") or getattr(project, "compliance", "")),
        "Duration": duration,
        "Generated At": datetime.now(ist).strftime("%Y-%m-%d %H:%M %Z"),
    }
    try:
        if getattr(project, "company", None):
            data["overview"]["Currency"] = getattr(project.company, "currency", "USD")
        else:
            data["overview"]["Currency"] = "USD"
    except Exception:
        data["overview"]["Currency"] = "USD"

    # Add discount to overview if present
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        data["overview"]["Discount"] = f"{discount_percentage}%"
        total_cost = sum(plan_entry.get("Cost", 0) for plan_entry in resourcing_plan)
        data["overview"]["Total Cost (After Discount)"] = f"${total_cost:,.2f}"

    data["activities"] = activities
    data["resourcing_plan"] = resourcing_plan

    # Keep discount_percentage in output for reference
    if discount_percentage and isinstance(discount_percentage, (int, float)) and discount_percentage > 0:
        data["discount_percentage"] = discount_percentage

    return data


async def generate_project_scope(db: AsyncSession, project) -> dict:
    """
    Generate project scope + architecture diagram + store architecture in DB + return combined JSON.
    """

    #  Ensure the project has a valid company reference (fallback to Sigmoid)
    if not getattr(project, "company_id", None):
        from app.utils import ratecards
        sigmoid = await ratecards.get_or_create_sigmoid_company(db)
        project.company_id = sigmoid.id
        await db.commit()
        await db.refresh(project)
        logger.info(f"Linked project {project.id} to Sigmoid company as fallback")

    tokenizer = tiktoken.get_encoding("cl100k_base")
    context_limit = 128000
    max_total_tokens = context_limit - 4000
    used_tokens = 0

    # ---------- Extract RFP ----------
    rfp_text = ""
    try:
        files: List[dict] = []
        if getattr(project, "files", None):
            try:
                files = [{"file_name": f.file_name, "file_path": f.file_path} for f in project.files]
            except Exception as e:
                logger.warning(f" Could not access project.files: {e}")
                files = []
        if files:
            rfp_text = await _extract_text_from_files(files)
    except Exception as e:
        logger.warning(f"File extraction for project {getattr(project, 'id', None)} failed: {e}")

    # ---------- Trim RFP text ----------
    rfp_tokens = tokenizer.encode(rfp_text or "")
    if len(rfp_tokens) > 5000:
        rfp_tokens = rfp_tokens[:5000]
    rfp_text = tokenizer.decode(rfp_tokens)
    used_tokens += len(rfp_tokens)

    # ---------- Retrieve KB context ----------
    fallback_fields = [
        getattr(project, "name", None),
        getattr(project, "domain", None),
        getattr(project, "complexity", None),
        getattr(project, "tech_stack", None),
        getattr(project, "use_cases", None),
        getattr(project, "compliance", None),
        str(getattr(project, "duration", "")) if getattr(project, "duration", None) else None,
    ]
    fallback_text = " ".join(f for f in fallback_fields if f and str(f).strip())

    # If completely empty, create a detailed specific prompt instead of returning empty scope
    if not (rfp_text.strip() or fallback_text.strip()):
        logger.warning(f"‚ö†Ô∏è No RFP text or project metadata for project {project.id}. Using detailed generic prompt.")
        fallback_text = """
Project Requirements:
- Project Type: Software Development Project
- Domain: Web Application Development
- Complexity: Medium
- Tech Stack: React, Node.js, PostgreSQL, AWS
- Duration: 6 months
- Team Size: 5-7 people

Project Scope:
Create a comprehensive project plan with the following phases:

1. Requirements & Planning Phase (1 month)
   - Gather and document requirements
   - Create technical specifications
   - Set up project infrastructure
   Owner: Project Manager
   Resources: Business Analyst, Technical Lead

2. Design Phase (1 month)
   - Design system architecture
   - Create UI/UX mockups
   - Database schema design
   Owner: Solution Architect
   Resources: UI/UX Designer, Database Administrator

3. Development Phase (2.5 months)
   - Frontend development (React)
   - Backend API development (Node.js)
   - Database implementation
   - Integration testing
   Owner: Technical Lead
   Resources: Frontend Developer, Backend Developer, QA Engineer

4. Testing & QA Phase (1 month)
   - Unit testing
   - Integration testing
   - User acceptance testing
   - Bug fixes
   Owner: QA Lead
   Resources: QA Engineer, Backend Developer

5. Deployment & Go-Live (0.5 months)
   - Production deployment
   - Performance optimization
   - Documentation
   - Training
   Owner: DevOps Engineer
   Resources: Technical Lead, Backend Developer

Generate activities with realistic start/end dates, proper role assignments, and meaningful descriptions.
"""

    kb_results = _rag_retrieve(rfp_text or fallback_text)
    kb_chunks = []
    stop = False
    for group in kb_results:
        for ch in group["chunks"]:
            chunk_tokens = len(tokenizer.encode(ch["content"]))
            if used_tokens + chunk_tokens > max_total_tokens:
                stop = True
                break
            kb_chunks.append(ch["content"])
            used_tokens += chunk_tokens
        if stop:
            break

    logger.info(
        f"Final RFP tokens: {len(rfp_tokens)}, KB tokens: {used_tokens - len(rfp_tokens)}, Total: {used_tokens}/{max_total_tokens}"
    )

    # ---------- Load questions.json (if exists) and build Q&A context ----------
    questions_context = None
    try:
        q_blob_name = f"{PROJECTS_BASE}/{project.id}/questions.json"
        if await azure_blob.blob_exists(q_blob_name):
            q_bytes = await azure_blob.download_bytes(q_blob_name)
            q_json = json.loads(q_bytes.decode("utf-8"))

            q_lines = []
            for category in q_json.get("questions", []):
                cat_name = category.get("category", "General")
                q_lines.append(f"### {cat_name}")
                for item in category.get("items", []):
                    q = item.get("question", "").strip()
                    a = item.get("user_understanding", "").strip() or "(unanswered)"
                    comment = item.get("comment", "").strip()
                    line = f"Q: {q}\nA: {a}"
                    if comment:
                        line += f"\nComment: {comment}"
                    q_lines.append(line)

            questions_context = "\n".join(q_lines)
            logger.info(f"Loaded {len(q_lines)} question lines for project {project.id}")
        else:
            logger.info(f"No questions.json found for project {project.id}, skipping Q&A context.")

    except Exception as e:
        logger.warning(f" Could not include questions.json context: {e}")
        questions_context = None

    # ---------- Fetch company rate card roles ----------
    rate_card_roles = []
    try:
        rate_map = await get_rate_map_for_project(db, project)
        rate_card_roles = list(rate_map.keys())
        logger.info(f"üìã Fetched {len(rate_card_roles)} roles from company rate card: {', '.join(rate_card_roles[:10])}")
    except Exception as e:
        logger.warning(f"‚ö†Ô∏è Could not fetch rate card roles: {e}")
        rate_card_roles = list(ROLE_RATE_MAP.keys())

    # ---------- Build + query ----------
    prompt = _build_scope_prompt(rfp_text, kb_chunks, project, questions_context=questions_context, rate_card_roles=rate_card_roles)
    try:
        # Step 1: Generate scope via Ollama with JSON format enforcement
        logger.info(f"ü§ñ Calling Ollama for scope generation... (prompt length: {len(prompt)} chars)")
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, format_json=True))
        logger.info(f"üìù Ollama raw response length: {len(raw_text)} chars")

        # Log more of the raw response to debug parsing issues
        logger.info(f"üìù Ollama response FIRST 1000 chars:\n{raw_text[:1000]}")
        logger.info(f"üìù Ollama response LAST 1000 chars:\n{raw_text[-1000:]}")

        if not raw_text or len(raw_text.strip()) < 50:
            logger.error(f"‚ùå Ollama returned empty or too short response: {len(raw_text)} chars")
            logger.error("   This usually means:")
            logger.error("   1. Ollama service is not running properly")
            logger.error("   2. The model (deepseek-r1) is not loaded")
            logger.error("   3. Out of memory or timeout")
            return {}

        raw = _extract_json(raw_text)

        # Log what was extracted
        logger.info(f"‚úÖ Extracted JSON keys: {list(raw.keys()) if isinstance(raw, dict) else 'NOT A DICT'}")
        if isinstance(raw, dict):
            logger.info(f"   - overview: {'present' if raw.get('overview') else 'MISSING'}")
            logger.info(f"   - activities: {len(raw.get('activities', []))} items")
            logger.info(f"   - resourcing_plan: {'present (will be auto-generated)' if 'resourcing_plan' in raw else 'not in raw'}")
            logger.info(f"   - project_summary: {'present' if raw.get('project_summary') else 'MISSING'}")
        else:
            logger.error(f"‚ùå Extracted result is not a dict! Type: {type(raw)}, Value: {raw}")
            return {}

        # Transform nested schema to flat schema if needed
        raw = _transform_nested_to_flat_schema(raw, project)

        # Log post-transformation
        logger.info(f"üìä After transformation - activities: {len(raw.get('activities', []))} items, "
                   f"overview: {'present' if raw.get('overview') else 'missing'}")

        # Validate that LLM actually generated content, not just structure
        if raw.get('activities'):
            activities = raw.get('activities', [])
            empty_fields_count = 0
            for act in activities:
                if (not act.get('Activities', '').strip() or
                    not act.get('Description', '').strip() or
                    act.get('Owner', '').lower() in ['unassigned', '']):
                    empty_fields_count += 1

            if empty_fields_count > len(activities) * 0.7:  # More than 70% are garbage
                logger.error(f"‚ùå LLM returned {empty_fields_count}/{len(activities)} activities with empty/invalid content!")
                logger.error("   This means Ollama generated JSON structure but NO actual content.")
                logger.error("   Check if:")
                logger.error("   1. Ollama service is running: curl http://localhost:11434/api/tags")
                logger.error("   2. Model is loaded: ollama list")
                logger.error("   3. Sufficient memory available")
                return {}
        else:
            logger.warning(f"‚ö†Ô∏è NO activities found in extracted JSON! This is a problem.")
            logger.warning(f"   Raw JSON structure: {json.dumps(raw, indent=2)[:500]}")

        cleaned_scope = await clean_scope(db, raw, project=project)
        # Update project fields from generated overview (just like finalize_scope)
        overview = cleaned_scope.get("overview", {})
        if overview:
            project.name = overview.get("Project Name") or project.name
            project.domain = overview.get("Domain") or project.domain
            project.complexity = overview.get("Complexity") or project.complexity
            project.tech_stack = overview.get("Tech Stack") or project.tech_stack
            project.use_cases = overview.get("Use Cases") or project.use_cases
            project.compliance = overview.get("Compliance") or project.compliance
            project.duration = str(overview.get("Duration") or project.duration)

            try:
                await db.commit()
                await db.refresh(project)
                logger.info(f" Project metadata updated from generated scope for project {project.id}")
            except Exception as e:
                logger.warning(f" Failed to update project metadata: {e}")


        # Step 2: Generate + store architecture diagram
        try:
            blob_base_path = f"{PROJECTS_BASE}/{getattr(project, 'id', 'unknown')}"
            db_file, arch_blob = await generate_architecture(
                db, project, rfp_text, kb_chunks, blob_base_path
            )
            cleaned_scope["architecture_diagram"] = arch_blob or None
        except Exception as e:
            logger.warning(f"Architecture diagram generation failed: {e}")
            cleaned_scope["architecture_diagram"] = None

        # Step 3: Auto-save finalized_scope.json in Azure Blob + DB
        try:
            from sqlalchemy import select
            result = await db.execute(
                select(models.ProjectFile).filter(
                    models.ProjectFile.project_id == project.id,
                    models.ProjectFile.file_name == "finalized_scope.json",
                )
            )
            old_file = result.scalars().first()
            if old_file:
                logger.info(f"Overwriting existing finalized_scope.json for project {project.id}")
            else:
                old_file = models.ProjectFile(
                    project_id=project.id,
                    file_name="finalized_scope.json",
                )

            blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"

            await azure_blob.upload_bytes(
                json.dumps(cleaned_scope, ensure_ascii=False, indent=2).encode("utf-8"),
                blob_name,
                overwrite=True, 
            )

            old_file.file_path = blob_name
            db.add(old_file)
            await db.commit()
            await db.refresh(old_file)

            logger.info(f" finalized_scope.json overwritten for project {project.id}")

        except Exception as e:
            logger.warning(f" Failed to auto-save finalized_scope.json: {e}")
        return cleaned_scope

    except Exception as e:
        logger.error(f"Ollama scope generation failed: {e}")
        return {}


async def regenerate_from_instructions(
    db: AsyncSession,
    project: models.Project,
    draft: dict,
    instructions: str
) -> dict:
    """
    Regenerate the project scope from user instructions using a creative AI-guided prompt.
    Enhances activity sequencing, roles, and effort estimates while preserving valid JSON structure.
    """
    logger.info(f" Regenerating scope for project {project.id} with creative AI response...")

    if not instructions or not instructions.strip():
        cleaned = await clean_scope(db, draft, project=project)
        return {**cleaned, "_finalized": True}


    prompt = f"""
You are an **expert AI project planner and delivery architect** responsible for maintaining a project scope in JSON format.

You are given:
1. The current draft project scope (JSON with keys: `overview`, `activities`, `resourcing_plan`).
2. The user‚Äôs latest change instructions.

Your task:
- **Understand** the user‚Äôs intent (instructions may be in natural language).
- **Regenerate** the scope accordingly:
  - Apply all user instructions faithfully.
  - Preserve structure and realism of the plan.
  - Re-calculate activity dates, dependencies, and efforts using the rules below.
  - Reflect improvements like ‚Äúoptimize‚Äù, ‚Äúsimplify‚Äù, ‚Äúrebalance‚Äù, or ‚Äúadd QA phase‚Äù.

---

### RULES OF MODIFICATION

####  Schema
- Preserve the same top-level keys: `overview`, `activities`, `resourcing_plan`.
- Every activity must have: "ID", "Activities", "Description", "Owner", "Resources",
- "Start Date", "End Date", "Effort Months"
- Use valid ISO dates (`yyyy-mm-dd`).
- Keep total duration ‚â§ 12 months.

**CRITICAL: What activities look like**
CORRECT activity example:
```json
{{
  "ID": 1,
  "Activities": "Project Initiation and Requirements Gathering",
  "Description": "Define project scope, gather requirements, create initial documentation",
  "Owner": "Project Manager",
  "Resources": "Business Analyst, Data Architect",
  "Start Date": "2025-01-15",
  "End Date": "2025-02-28",
  "Effort Months": 1.5
}}
```

WRONG activity example (DO NOT DO THIS):
```json
{{
  "ID": 1,
  "Activities": "Project Manager",  ‚Üê WRONG! This is a role name, not an activity!
  "Description": "",  ‚Üê WRONG! Must have meaningful description!
  "Owner": "Unassigned",  ‚Üê WRONG! Must have a real owner!
  "Resources": "",
  "Start Date": "2025-01-15",
  "End Date": "2025-02-15",
  "Effort Months": 1
}}
```

####  Temporal Adjustment Rules
Use these to keep the schedule consistent and continuous.

**Add new activity (bottom)**  
- Append at the end.  
- Start date = 10 days *before* the current latest end_date.  
- End date = start_date + duration derived from effort_days.  
- Allow small overlap (10-15 %) with the last activity to maximize parallelism.

**Add new activity (in middle)**  
- Insert between the target activities without disturbing global schedule.  
- Preceding activity‚Äôs end date remains fixed.  
- Following activity‚Äôs start shifts minimally to maintain continuity.  
- Only local dates adjust; efforts remain unchanged for following activities.

**Delete activity**  
- Remove it completely.  
- Do not introduce gaps; subsequent activities retain start/end dates.

**Split activity into two**  
- Divide one activity into two consecutive ones.  
- Combined effort_days = original.  
- Combined duration = original.  
- Other activities‚Äô dates stay the same.

**Merge two activities**
- Combine both into one.
- start_date = min(start of both)
- end_date = max(end of both)
- effort_days = sum(efforts of both)

####  Role Management Rules
Critical: When user requests to add or remove roles, you MUST update BOTH activities and resourcing_plan.

**IMPORTANT: All changes are INCREMENTAL - preserve existing activities unless explicitly deleted!**

**Remove a role (e.g., "remove Business Analyst")**:
1. Keep ALL existing activities
2. Find all activities where the role is the Owner
3. Reassign those activities to another appropriate role
4. Remove the role from ALL Resources fields across all activities
5. DO NOT delete any activities - only change role assignments
6. Example: If removing "Business Analyst":
   - Activity: "Owner": "Business Analyst" ‚Üí change to "Owner": "Product Manager"
   - Activity: "Resources": "Business Analyst, Data Engineer" ‚Üí change to "Resources": "Data Engineer"
   - Keep ALL other activities unchanged
   - resourcing_plan: will be auto-calculated

**Add more of an existing role (e.g., "add 1 more Backend Developer")**:
1. **CRITICAL**: Keep ALL existing activities and roles
2. "Add 1 more" means INCREASE allocation, not replace
3. To increase Backend Developer allocation:
   - Add "Backend Developer" to Resources field of MORE existing activities
   - OR extend date ranges of activities that already have Backend Developer
   - OR create 1-2 NEW activities specifically for Backend Developer
4. **DO NOT remove any existing activities or roles**
5. Example: If you have 10 activities and "add 1 Backend Developer":
   - Original: 10 activities with Backend Developer in 3 of them
   - After: Same 10 activities PLUS Backend Developer added to 2-3 more activities
   - Result: Backend Developer effort increases from 3 months to 5-6 months

**Add a new role type (e.g., "add Security Engineer")**:
1. **CRITICAL**: Keep ALL existing activities and roles
2. Add new activities for this role OR add to Resources field of existing activities
3. DO NOT remove any existing activities
4. The resourcing_plan will be auto-generated based on activities

####  Discount Rules
When user requests a discount (e.g., "apply 5% discount", "give 10% discount"):
1. **DO NOT change activities, dates, or efforts**
2. **ONLY note the discount percentage in a special field**
3. Add a new field: "discount_percentage": <number> (e.g., 5 for 5%, 10 for 10%)
4. The discount will be applied automatically during cost calculation
5. Keep all activities, roles, and resourcing_plan calculations unchanged

### Scheduling Rules
- Activities should follow **semi-parallel execution** ‚Äî overlap realistically but maintain logical order.
- If two activities are **independent**, overlap their timelines by **70‚Äì80%** of their duration (not full overlap).
- If one activity **depends** on another, allow a small overlap of **10-15%** near the end of the predecessor if feasible.
- Avoid full serialization unless strictly required by dependency.
- Avoid full parallelism where all tasks start together ‚Äî stagger independent ones by **10-15%**.
- Ensure overall project duration stays **‚â§ 12 months**.
- The first activity must always start today.
---

### Regeneration Logic
- Clean and re-order activities logically.
- Maintain coherent dependencies and sequential flow.
- Adjust `overview.duration_months` automatically based on new total project span.
- Keep resource roles realistic and consistent with activities (Backend Developer, Data Engineer, QA Analyst, etc.).
- Reflect optimization or simplification requests (e.g., reduce redundant steps, consolidate phases).

---

###  Output Rules
- Output **only valid JSON** ‚Äî no markdown, no explanations, no reasoning.
- Must include:
  - `overview` ‚Üí Project metadata (name, domain, complexity, tech stack, etc.)
  - `activities` ‚Üí COMPLETE updated list with ALL modifications applied
  - `resourcing_plan` ‚Üí OPTIONAL (will be auto-calculated from activities)
  - `discount_percentage` ‚Üí OPTIONAL (only if user requested discount, e.g., 5 for 5%, 10 for 10%)
- **CRITICAL**: If user says "remove [role]", that role MUST NOT appear in ANY activity's Owner or Resources field
- **CRITICAL**: If user says "add 1 more [role]", ADD to existing activities, DO NOT replace them
- **CRITICAL**: If user says "apply X% discount", include "discount_percentage": X in output
- **Dont change schema or field names.**
- **PRESERVE all activities** - only modify/add/remove specific items mentioned by user

---

User Instructions:
{instructions}

Current Draft Scope:
{json.dumps(draft, indent=2, ensure_ascii=False)}

Return only the updated JSON.
"""


    # ---- Query Ollama creatively with JSON enforcement ----
    # Use lower temperature for more consistent instruction-following
    try:
        raw_text = await anyio.to_thread.run_sync(lambda: ollama_chat(prompt, temperature=0.2, format_json=True))
        logger.info(f"ü§ñ LLM response length: {len(raw_text)} chars")
        logger.debug(f"LLM raw response (first 500 chars): {raw_text[:500]}")
        updated_scope = _extract_json(raw_text)

        logger.info(f"üìä Extracted scope structure: overview={bool(updated_scope.get('overview'))}, "
                   f"activities={len(updated_scope.get('activities', []))}, "
                   f"resourcing_plan={len(updated_scope.get('resourcing_plan', []))}")

        # Validate activity count - prevent accidental scope replacement
        original_activity_count = len(draft.get('activities', []))
        new_activity_count = len(updated_scope.get('activities', []))
        is_removal_instruction = any(word in instructions.lower() for word in ['remove', 'delete'])

        # Advanced validation: Check if activities are valid/meaningful
        activities_are_valid = True
        validation_failures = []

        if updated_scope.get('activities'):
            unassigned_count = sum(1 for act in updated_scope['activities'] if act.get('Owner', '').lower() in ['unassigned', ''])
            empty_desc_count = sum(1 for act in updated_scope['activities'] if not act.get('Description', '').strip())

            # Check if activity names are just role names (common LLM mistake)
            common_roles = ['project manager', 'business analyst', 'data architect', 'data engineer',
                           'backend developer', 'frontend developer', 'qa engineer', 'devops engineer',
                           'cloud architect', 'data analyst', 'ux designer']
            role_name_activities = sum(1 for act in updated_scope['activities']
                                      if act.get('Activities', '').lower().strip() in common_roles)

            # Check if all activities have identical dates (suspicious)
            dates = [(act.get('Start Date'), act.get('End Date')) for act in updated_scope['activities']]
            unique_date_pairs = len(set(dates))

            # Validation thresholds
            if unassigned_count > new_activity_count * 0.5:  # More than 50% unassigned
                activities_are_valid = False
                validation_failures.append(f"{unassigned_count}/{new_activity_count} activities have Unassigned owner")

            if empty_desc_count > new_activity_count * 0.5:  # More than 50% empty descriptions
                activities_are_valid = False
                validation_failures.append(f"{empty_desc_count}/{new_activity_count} activities have empty descriptions")

            if role_name_activities > new_activity_count * 0.3:  # More than 30% are just role names
                activities_are_valid = False
                validation_failures.append(f"{role_name_activities}/{new_activity_count} activities are named after roles (e.g. 'Project Manager', 'Data Engineer')")

            if unique_date_pairs == 1 and new_activity_count > 1:  # All activities have same dates
                activities_are_valid = False
                validation_failures.append(f"All {new_activity_count} activities have identical dates: {dates[0]}")

        # If LLM significantly reduced activities OR created invalid activities, restore original
        if (new_activity_count < (original_activity_count * 0.7) and not is_removal_instruction) or not activities_are_valid:
            if not activities_are_valid:
                logger.error(f"‚ùå LLM GENERATED INVALID ACTIVITIES!")
                for failure in validation_failures:
                    logger.error(f"   - {failure}")
            else:
                logger.error(f"‚ùå LLM LOST TOO MANY ACTIVITIES! Original: {original_activity_count}, New: {new_activity_count}")

            logger.error(f"   User instruction: '{instructions[:100]}'")
            logger.error(f"   üîß Auto-restoring original activities to prevent data loss")

            # Restore original activities
            updated_scope["activities"] = draft.get("activities", [])
            if "resourcing_plan" not in updated_scope or not updated_scope.get("resourcing_plan"):
                updated_scope["resourcing_plan"] = draft.get("resourcing_plan", [])

            logger.info(f"   ‚úÖ Restored {len(updated_scope['activities'])} valid activities from draft")

        # Log roles found in activities
        if updated_scope.get('activities'):
            owners = set(act.get('Owner', 'Unknown') for act in updated_scope['activities'])
            all_resources = set()
            for act in updated_scope['activities']:
                resources = act.get('Resources', '')
                if resources:
                    all_resources.update(r.strip() for r in str(resources).split(',') if r.strip())
            all_roles = owners | all_resources
            logger.info(f"üé≠ Roles in LLM response - Owners: {owners}, Resources: {all_resources}")

            # Validate that "remove" instructions were followed
            if instructions and 'remove' in instructions.lower():
                for role in all_roles:
                    if role.lower() in instructions.lower() and 'remove' in instructions.lower():
                        logger.error(f"‚ùå LLM FAILED to remove '{role}' - still present in activities despite user instruction!")

            # Validate that "add" instructions were followed
            if instructions and 'add' in instructions.lower():
                # This is harder to validate automatically, but we log for manual inspection
                logger.info(f"‚ÑπÔ∏è User requested to add role(s). Current roles: {all_roles}")

        # Post-processing fallback: manually remove roles if LLM failed
        if instructions and 'remove' in instructions.lower() and updated_scope.get('activities'):
            # Extract role to remove from instructions (basic pattern matching)
            import re
            # Pattern to match "remove <role>" where role can be multi-word
            # Matches everything after "remove" until end of string or common delimiters
            remove_pattern = r'remove\s+([a-zA-Z\s]+?)(?:\s*(?:from|,|\.|\band\b|$))'
            match = re.search(remove_pattern, instructions.lower(), re.IGNORECASE)
            if match:
                role_to_remove = match.group(1).strip()
                logger.info(f"üîß Post-processing: attempting to remove '{role_to_remove}'")

                # Track if we made changes
                changes_made = False

                # Process each activity
                for act in updated_scope['activities']:
                    # Check if this role is the owner
                    if act.get('Owner', '').lower() == role_to_remove or role_to_remove in act.get('Owner', '').lower():
                        # Find a replacement owner from resources or use a default
                        resources = act.get('Resources', '')
                        if resources and resources.strip():
                            # Use the first resource as the new owner
                            new_owner = resources.split(',')[0].strip()
                            # Remove new owner from resources to avoid duplication
                            remaining_resources = [r.strip() for r in resources.split(',')[1:] if r.strip()]
                            act['Owner'] = new_owner
                            act['Resources'] = ', '.join(remaining_resources)
                            logger.info(f"  ‚Üí Reassigned activity '{act.get('Activities', 'Unknown')}' from removed role to '{new_owner}'")
                            changes_made = True
                        else:
                            # No resources available, use a generic default
                            act['Owner'] = 'Project Manager'
                            logger.info(f"  ‚Üí Reassigned activity '{act.get('Activities', 'Unknown')}' from removed role to 'Project Manager'")
                            changes_made = True

                    # Remove from resources field
                    if act.get('Resources'):
                        resources_list = [r.strip() for r in str(act['Resources']).split(',') if r.strip()]
                        # Filter out the role to remove (case-insensitive partial match)
                        filtered_resources = [r for r in resources_list
                                             if role_to_remove not in r.lower() and r.lower() != role_to_remove]
                        if len(filtered_resources) != len(resources_list):
                            act['Resources'] = ', '.join(filtered_resources)
                            changes_made = True

                if changes_made:
                    logger.info(f"‚úÖ Post-processing successfully removed role '{role_to_remove}' from activities")

        # Post-processing: parse discount percentage from instructions
        if instructions:
            import re
            # Pattern to match discount requests: "5% discount", "apply 10% discount", "give 15% discount", etc.
            discount_patterns = [
                r'(\d+)\s*%\s*discount',
                r'discount\s+(?:of\s+)?(\d+)\s*%',
                r'apply\s+(\d+)\s*%',
                r'give\s+(\d+)\s*%',
            ]
            discount_found = False
            for pattern in discount_patterns:
                match = re.search(pattern, instructions.lower())
                if match:
                    discount_percentage = int(match.group(1))
                    logger.info(f"üí∞ Post-processing: detected {discount_percentage}% discount request")

                    # Add discount to updated_scope if not already present
                    if "discount_percentage" not in updated_scope or not updated_scope.get("discount_percentage"):
                        updated_scope["discount_percentage"] = discount_percentage
                        logger.info(f"  ‚Üí Added discount_percentage: {discount_percentage}")
                    discount_found = True
                    break

            if not discount_found and any(word in instructions.lower() for word in ['discount', 'reduction', 'reduce cost']):
                logger.warning(f"‚ö†Ô∏è User mentioned discount but couldn't parse percentage. Instructions: {instructions[:100]}")

        # Safety check: if LLM returned empty activities, preserve original
        if not updated_scope.get("activities") or len(updated_scope.get("activities", [])) == 0:
            logger.warning(f"‚ö†Ô∏è LLM returned empty activities - preserving original draft activities")
            logger.info(f"üìã Original draft had {len(draft.get('activities', []))} activities")
            # Preserve original activities and resourcing_plan, but update overview if changed
            updated_scope["activities"] = draft.get("activities", [])
            if "resourcing_plan" not in updated_scope or not updated_scope.get("resourcing_plan"):
                updated_scope["resourcing_plan"] = draft.get("resourcing_plan", [])

        cleaned = await clean_scope(db, updated_scope, project=project)
        logger.info(f"‚úÖ Cleaned scope: {len(cleaned.get('activities', []))} activities, "
                   f"{len(cleaned.get('resourcing_plan', []))} resources")

    except Exception as e:
        logger.error(f" Creative regeneration failed: {e}")
        cleaned = await clean_scope(db, draft, project=project)

    # ---- Update project metadata from overview ----
    overview = cleaned.get("overview", {})
    if overview:
        project.name = overview.get("Project Name") or project.name
        project.domain = overview.get("Domain") or project.domain
        project.complexity = overview.get("Complexity") or project.complexity
        project.tech_stack = overview.get("Tech Stack") or project.tech_stack
        project.use_cases = overview.get("Use Cases") or project.use_cases
        project.compliance = overview.get("Compliance") or project.compliance
        project.duration = str(overview.get("Duration") or project.duration)
        await db.commit()
        await db.refresh(project)
        logger.info(f" Project metadata synced for project {project.id}")

    # ---- Overwrite finalized_scope.json in Blob ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project.id,
            models.ProjectFile.file_name == "finalized_scope.json",
        )
    )
    old_file = result.scalars().first() or models.ProjectFile(
        project_id=project.id, file_name="finalized_scope.json"
    )

    blob_name = f"{PROJECTS_BASE}/{project.id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(cleaned, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )
    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f" Creative finalized_scope.json regenerated for project {project.id}")
    return {**cleaned, "_finalized": True}


async def finalize_scope(
    db: AsyncSession,
    project_id: str,
    scope_data: dict
) -> tuple[models.ProjectFile, dict]:
    """
    Finalize the project scope without LLM ‚Äî just clean, validate sequencing,
    update metadata, and save finalized_scope.json.
    """

    logger.info(f"Finalizing scope (no LLM) for project {project_id}...")

    # ---- Load project ----
    result = await db.execute(
        select(models.Project)
        .options(selectinload(models.Project.company))
        .filter(models.Project.id == project_id)
    )
    project = result.scalars().first()
    if not project:
        raise ValueError(f"Project {project_id} not found")

    # ---- Step 1: Clean draft ----
    finalized = await clean_scope(db, scope_data, project=project)
    overview = finalized.get("overview", {})

    # ---- Step 2: Update project metadata ----
    if overview:
        project.name = overview.get("Project Name") or project.name
        project.domain = overview.get("Domain") or project.domain
        project.complexity = overview.get("Complexity") or project.complexity
        project.tech_stack = overview.get("Tech Stack") or project.tech_stack
        project.use_cases = overview.get("Use Cases") or project.use_cases
        project.compliance = overview.get("Compliance") or project.compliance
        project.duration = str(overview.get("Duration") or project.duration)
        await db.commit()
        await db.refresh(project)

    # ---- Step 3: Save finalized_scope.json ----
    result = await db.execute(
        select(models.ProjectFile).filter(
            models.ProjectFile.project_id == project_id,
            models.ProjectFile.file_name == "finalized_scope.json"
        )
    )
    old_file = result.scalars().first()
    if old_file:
        logger.info(f" Overwriting existing finalized_scope.json for project {project_id}")
    else:
        old_file = models.ProjectFile(
            project_id=project_id,
            file_name="finalized_scope.json",
        )

    blob_name = f"{PROJECTS_BASE}/{project_id}/finalized_scope.json"
    await azure_blob.upload_bytes(
        json.dumps(finalized, ensure_ascii=False, indent=2).encode("utf-8"),
        blob_name,
        overwrite=True,
    )

    old_file.file_path = blob_name
    db.add(old_file)
    await db.commit()
    await db.refresh(old_file)

    logger.info(f" Finalized scope saved (no LLM) for project {project_id}")
    return old_file, {**finalized, "_finalized": True}