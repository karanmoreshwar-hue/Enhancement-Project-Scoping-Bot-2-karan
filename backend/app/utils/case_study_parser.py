"""
Case Study Parser - Extract structured case study data from PPT files.

This module parses PowerPoint files containing case studies and extracts
structured information including: Client Name, Overview, Solution, and Impact.
"""

import re
import logging
from typing import Dict, List, Optional
from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_case_study_from_ppt(file_path: str) -> List[Dict[str, str]]:
    """
    Parse a PPT file and extract structured case study information.

    Expected PPT structure:
    - Each case study should have slides with recognizable headers
    - Headers can be: "Client", "Overview", "Solution", "Impact", "Results", etc.
    - The parser will try to intelligently extract these fields

    Args:
        file_path: Path to the PPT file

    Returns:
        List of case study dictionaries, each containing:
        {
            "client_name": str,
            "overview": str,
            "solution": str,
            "impact": str,
            "full_text": str  # Combined text for embedding
        }
    """
    try:
        prs = Presentation(file_path)
        case_studies = []
        current_case_study = None

        for slide_idx, slide in enumerate(prs.slides):
            slide_text = _extract_slide_text(slide)

            # Check if this is a new case study (title slide)
            if _is_case_study_title_slide(slide_text):
                # Save previous case study if exists
                if current_case_study and _is_valid_case_study(current_case_study):
                    current_case_study["full_text"] = _build_full_text(current_case_study)
                    case_studies.append(current_case_study)
                    logger.info(f"📚 Extracted case study: {current_case_study.get('client_name', 'Unknown')}")

                # Start new case study
                current_case_study = {
                    "client_name": _extract_client_name(slide_text),
                    "overview": "",
                    "solution": "",
                    "impact": "",
                    "slide_range": f"{slide_idx + 1}"
                }

            # Extract fields from current slide
            if current_case_study:
                _extract_case_study_fields(slide_text, current_case_study)
                # Update slide range
                start_slide = current_case_study["slide_range"].split("-")[0]
                current_case_study["slide_range"] = f"{start_slide}-{slide_idx + 1}"

        # Don't forget the last case study
        if current_case_study and _is_valid_case_study(current_case_study):
            current_case_study["full_text"] = _build_full_text(current_case_study)
            case_studies.append(current_case_study)
            logger.info(f"📚 Extracted case study: {current_case_study.get('client_name', 'Unknown')}")

        # Post-process: Extract overview from long client names if overview is missing
        for case_study in case_studies:
            _extract_overview_from_client_name(case_study)

        logger.info(f"✅ Parsed {len(case_studies)} case studies from PPT")
        return case_studies

    except Exception as e:
        logger.error(f"❌ Failed to parse case study PPT: {e}")
        return []


def _extract_slide_text(slide) -> str:
    """Extract all text from a slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_parts.append(shape.text.strip())
    return "\n".join(text_parts)


def _is_case_study_title_slide(text: str) -> bool:
    """
    Determine if this slide is the start of a new case study.

    Indicators:
    - Contains "case study" in title
    - Contains "client" or "customer" keywords
    - Starts with a company/organization name pattern
    """
    text_lower = text.lower()

    # Check for explicit case study indicators
    if any(keyword in text_lower for keyword in ["case study", "client story", "customer success"]):
        return True

    # Check for client/customer name patterns
    if re.search(r"(client|customer):\s*\w+", text_lower):
        return True

    # Check if first line looks like a company name (capitalized, short)
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    if lines and len(lines[0]) < 100 and lines[0][0].isupper():
        # Likely a title slide with company name
        return True

    return False


def _extract_client_name(text: str) -> str:
    """Extract client/company name from slide text."""
    lines = [l.strip() for l in text.split("\n") if l.strip()]

    # Look for explicit "Client:" or "Customer:" labels
    for line in lines:
        match = re.search(r"(?:client|customer|company):\s*(.+)", line, re.IGNORECASE)
        if match:
            return match.group(1).strip()

    # Look for "case study" pattern
    for line in lines:
        match = re.search(r"(.+?)\s+case\s+study", line, re.IGNORECASE)
        if match:
            return match.group(1).strip()

    # Fallback: use first non-empty line as client name
    # If first line is very long (>150 chars), it likely contains description too
    # In that case, extract just the company name part before dash or parenthesis
    if lines:
        first_line = lines[0]

        # If line is long and contains dash/description, split it
        if len(first_line) > 150:
            # Try to extract just the company name before en dash (–) or em dash (—)
            # Note: Don't include regular hyphen (-) as it's used in names like "SK-II"
            match = re.match(r"^([^–—]+?)(?:\s*[–—]\s*)", first_line)
            if match:
                # Return just the company name part
                return match.group(1).strip()

        return first_line

    return "Unknown Client"


def _extract_case_study_fields(text: str, case_study: Dict[str, str]) -> None:
    """
    Extract overview, solution, and impact from slide text.

    Looks for section headers and extracts content below them.
    """
    text_lower = text.lower()
    lines = text.split("\n")

    # Overview/Background/Challenge
    if any(keyword in text_lower for keyword in ["overview", "background", "challenge", "problem", "about"]):
        content = _extract_section_content(text, ["overview", "background", "challenge", "problem", "about"])
        if content and len(content) > len(case_study["overview"]):
            case_study["overview"] = content

    # Solution/Approach/Implementation
    if any(keyword in text_lower for keyword in ["solution", "approach", "implementation", "methodology", "how we"]):
        content = _extract_section_content(text, ["solution", "approach", "implementation", "methodology", "how we"])
        if content and len(content) > len(case_study["solution"]):
            case_study["solution"] = content

    # Impact/Results/Outcomes/Benefits
    if any(keyword in text_lower for keyword in ["impact", "result", "outcome", "benefit", "achievement", "success"]):
        content = _extract_section_content(text, ["impact", "result", "outcome", "benefit", "achievement", "success"])
        if content and len(content) > len(case_study["impact"]):
            case_study["impact"] = content


def _extract_section_content(text: str, keywords: List[str]) -> str:
    """Extract content under a section header matching any of the keywords."""
    lines = text.split("\n")
    content_lines = []
    capturing = False

    for line in lines:
        line_lower = line.lower().strip()

        # Check if this line is a header matching our keywords
        if any(keyword in line_lower for keyword in keywords):
            capturing = True
            # Skip the header line itself
            continue

        # Check if we hit another section header (stop capturing)
        if capturing and line_lower and line[0].isupper() and len(line.strip()) < 50:
            # Likely another section header
            if any(stop_word in line_lower for stop_word in ["client", "company", "overview", "solution", "impact", "result", "challenge"]):
                break

        # Capture content
        if capturing and line.strip():
            content_lines.append(line.strip())

    return " ".join(content_lines).strip()


def _is_valid_case_study(case_study: Dict[str, str]) -> bool:
    """Check if case study has minimum required fields."""
    # Must have at least client name and one of the key fields
    if not case_study.get("client_name") or case_study.get("client_name") == "Unknown Client":
        return False

    # Must have at least one substantive field
    has_content = any([
        len(case_study.get("overview", "")) > 20,
        len(case_study.get("solution", "")) > 20,
        len(case_study.get("impact", "")) > 20
    ])

    return has_content


def _extract_overview_from_client_name(case_study: Dict[str, str]) -> None:
    """
    If client_name is very long and overview is empty, extract the description
    part from client_name and use it as overview.

    Example:
    Input:  client_name = "SK-II (P&G) – Global prestige skincare brand..."
    Output: client_name = "SK-II (P&G)"
            overview = "Global prestige skincare brand..."
    """
    client_name = case_study.get("client_name", "")
    overview = case_study.get("overview", "")

    # Only process if overview is empty and client_name is long
    if overview or len(client_name) < 100:
        return

    # Look for en dash (–) or em dash (—) separators
    # These often separate company name from description
    # Note: Don't include regular hyphen (-) as it's used in names like "SK-II"
    match = re.match(r"^(.+?)\s*[–—]\s*(.+)$", client_name)
    if match:
        company_part = match.group(1).strip()
        description_part = match.group(2).strip()

        # Only split if description part is substantial (>50 chars)
        if len(description_part) > 50:
            case_study["client_name"] = company_part
            case_study["overview"] = description_part
            logger.info(f"✂️ Split long client name into client + overview for: {company_part}")
            return

    # Alternative: Look for parentheses with long content after
    match = re.match(r"^(.+?)\s*\(([^)]+)\)\s*[–—]?\s*(.+)$", client_name)
    if match:
        company_part = f"{match.group(1).strip()} ({match.group(2).strip()})"
        description_part = match.group(3).strip()

        if len(description_part) > 50:
            case_study["client_name"] = company_part
            case_study["overview"] = description_part
            logger.info(f"✂️ Split long client name into client + overview for: {company_part}")


def _build_full_text(case_study: Dict[str, str]) -> str:
    """Build a full text representation for embedding generation."""
    parts = [
        f"Client: {case_study.get('client_name', '')}",
        f"Overview: {case_study.get('overview', '')}",
        f"Solution: {case_study.get('solution', '')}",
        f"Impact: {case_study.get('impact', '')}"
    ]
    return "\n\n".join(parts).strip()


def extract_all_text_from_ppt(file_path: str) -> str:
    """
    Fallback: Extract all text from PPT as a single string.
    Used when structured parsing fails.
    """
    try:
        prs = Presentation(file_path)
        all_text = []

        for slide in prs.slides:
            slide_text = _extract_slide_text(slide)
            if slide_text:
                all_text.append(slide_text)

        return "\n\n".join(all_text).strip()

    except Exception as e:
        logger.error(f"❌ Failed to extract text from PPT: {e}")
        return ""