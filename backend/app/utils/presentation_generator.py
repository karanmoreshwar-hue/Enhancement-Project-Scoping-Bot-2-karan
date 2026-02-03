from __future__ import annotations
import io
import json
import logging
import asyncio
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.utils.ai_clients import get_async_azure_client, AZURE_OPENAI_DEPLOYMENT
from app.utils.export import THEME
from app.utils import azure_blob
import os
import subprocess
import tempfile

logger = logging.getLogger(__name__)

TEMPLATE_PROMPT = """
You are an expert Presentation Designer & Strategist. Your goal is to create a compelling Pitch Deck script based on the provided Project Scope.
The output must be a valid JSON array of slide objects.

Project Scope:
{scope_json}

INSTRUCTIONS:
1. Create a 10-12 slide pitch deck that tells a story: Problem -> Solution -> Technical Approach -> Value -> Timeline -> Team.
2. Each slide must have:
   - "title": Short, punchy title.
   - "content": Bullet points (max 5) summarizing the key message.
   - "notes": Speaker notes for the presenter.
3. Be professional, persuasive, and concise. Use active voice.
4. Do NOT simply derive the content from the scope directly. Synthesize it into a pitch.
5. For the "Team" slide, use placeholders if names aren't in the scope.
6. For "Timeline", summarize key phases.
7. Return start and end with ```json and ```.

Output Format:
[
  {{
    "title": "Slide Title",
    "content": ["Point 1", "Point 2"],
    "notes": "Speaker notes..."
  }}
]
"""

async def generate_smart_pptx(scope: Dict[str, Any]) -> io.BytesIO:
    """
    Generates a smart AI-written pitch deck from the project scope.
    """
    try:
        # 1. Generate Script using LLM
        logger.info("🧠 Generating pitch deck script using LLM...")
        script = await _generate_script_with_llm(scope)
        
        # 2. Fetch Architecture Diagram (if any)
        diagram_data = await _fetch_architecture_diagram(scope)

        # 3. Render to PPTX
        logger.info(f"🎨 Rendering {len(script)} slides to PPTX...")
        pptx_buffer = _render_pptx(script, scope, diagram_data)
        
        return pptx_buffer

    except Exception as e:
        logger.error(f"❌ Smart Pitch Deck generation failed: {e}")
        raise

async def _generate_script_with_llm(scope: Dict[str, Any]) -> List[Dict[str, Any]]:
    client = get_async_azure_client()
    
    # Minimize scope to reduce token usage if needed, but for now send vital parts
    # Filter out massive lists if they exist to avoid token limits, keeping summaries
    scope_summary = {
        "overview": scope.get("overview"),
        "project_summary": scope.get("project_summary"),
        "activities": scope.get("activities", [])[:10], # First 10 activities for context
        "resourcing_plan": scope.get("resourcing_plan", [])[:5] # Summary of resources
    }
    
    prompt = TEMPLATE_PROMPT.format(scope_json=json.dumps(scope_summary, indent=2))

    response = await client.chat.completions.create(
        model=AZURE_OPENAI_DEPLOYMENT,
        messages=[
            {"role": "system", "content": "You are a helpful AI assistant that generates JSON for presentations."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.7,
        max_tokens=2000
    )

    content = response.choices[0].message.content
    
    # Extract JSON from markdown code block if present
    if "```json" in content:
        content = content.split("```json")[1].split("```")[0].strip()
    elif "```" in content:
        content = content.split("```")[1].strip()
        
    try:
        return json.loads(content)
    except json.JSONDecodeError:
        logger.error(f"Failed to decode LLM response: {content}")
        # Fallback to a basic structure if LLM fails
        return [
            {
                "title": "Project Scope Overview",
                "content": ["Failed to generate full script.", "Please refer to standard export."],
                "notes": "AI generation error."
            }
        ]

async def _fetch_architecture_diagram(scope: Dict[str, Any]) -> bytes | None:
    """
    Fetches the architecture diagram from blob storage.
    Handles PNG download and SVG->PNG conversion.
    """
    arch_path = scope.get("architecture_diagram")
    if not arch_path:
        return None

    logger.info(f"🔍 Fetching architecture diagram for PPTX: {arch_path}")
    img_bytes = None
    is_svg = arch_path.lower().endswith(".svg")

    try:
        # Try downloading the image
        try:
            img_bytes = await azure_blob.download_bytes(arch_path, timeout=15)
        except Exception:
            # If specific blob not found, pass to fallback check
            pass
        
        # If failed or empty, check for fallback
        if not img_bytes and not is_svg:
             if arch_path.endswith(".png"):
                svg_path = arch_path[:-4] + ".svg"
                if await azure_blob.blob_exists(svg_path):
                    is_svg = True
                    # Update path to use the fallback found
                    arch_path = svg_path
        
        if is_svg:
            try:
                # Download the SVG if we haven't already
                if not img_bytes:
                    img_bytes = await azure_blob.download_bytes(arch_path, timeout=15)

                if img_bytes:
                    # Create temp files for conversion
                    with tempfile.NamedTemporaryFile(suffix=".svg", delete=False) as f_in:
                        f_in.write(img_bytes)
                        temp_svg = f_in.name
                    
                    temp_png = temp_svg.replace(".svg", ".png")
                    
                    try:
                        # Find conversion script (assuming it's in backend root)
                        # Current file is in backend/app/utils
                        backend_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                        script_path = os.path.join(backend_root, "convert_svg.js")
                        
                        if os.path.exists(script_path):
                            logger.info(f"🔄 Converting SVG to PNG: {temp_svg} -> {temp_png}")
                            subprocess.run(
                                ["node", script_path, temp_svg, temp_png],
                                check=True,
                                capture_output=True
                            )
                            
                            # Read back the PNG
                            if os.path.exists(temp_png):
                                with open(temp_png, "rb") as f_png:
                                    img_bytes = f_png.read()
                                logger.info(f"✅ SVG conversion successful: {len(img_bytes)} bytes")
                                return img_bytes
                            else:
                                logger.error("❌ SVG conversion failed: Output PNG not found")
                        else:
                            logger.error(f"❌ Conversion script not found at {script_path}")
                            
                    except subprocess.CalledProcessError as e:
                        logger.error(f"❌ SVG conversion failed (node error): {e.stderr.decode()}")
                    except Exception as e:
                        logger.error(f"❌ SVG conversion failed: {e}")
                    finally:
                        # Cleanup
                        if os.path.exists(temp_svg):
                            os.unlink(temp_svg)
                        if os.path.exists(temp_png):
                            os.unlink(temp_png)

            except Exception as e:
                logger.error(f"Failed handling SVG conversion: {e}")
            
            # If we reached here with is_svg=True and no bytes returned, it failed.
            return None

        return img_bytes

    except Exception as e:
        logger.warning(f"Failed to fetch architecture diagram: {e}")
        return None

def _render_pptx(script: List[Dict[str, Any]], scope: Dict[str, Any], diagram_bytes: bytes | None = None) -> io.BytesIO:
    prs = Presentation()
    
    # Use theme colors
    header_bg = THEME["header_bg"] # e.g. #BDD7EE
    
    for slide_data in script:
        layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(layout)
        
        # INSERT DIAGRAM AFTER SLIDE 3 (Just after Problem/Solution typically)
        # OR if explicitly requested in script. For now, let's just insert it at the end or if we see a "Technical" slide.
        # simpler: Insert it as a dedicated slide after the first 2 slides (Introduction)
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get("title", "Untitled Slide")
        
        # Content (Bullets)
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        
        content_items = slide_data.get("content", [])
        if isinstance(content_items, str):
            content_items = [content_items]
            
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            
        # Speaker Notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data.get("notes")

    # --- Insert Architecture Diagram Slide ---
    if diagram_bytes:
        try:
            logger.info("🖼️ Adding Architecture Diagram slide...")
            slide_layout = prs.slide_layouts[5] # Title Only
            
            # Insert at index 3 (4th slide) if possible, else append
            target_index = min(3, len(prs.slides))
            # python-pptx doesn't support insert_slide easily on the slides collection directly in older versions, 
            # but usually we just append. Let's append to keep it safe.
            slide = prs.slides.add_slide(slide_layout)
            
            # Move slide to index 3? pptx is tricky with reordering.
            # Let's just keep it at the end or append it. 
            # Actually, let's append it right after the script.
            
            title = slide.shapes.title
            title.text = "System Architecture"
            
            img_stream = io.BytesIO(diagram_bytes)
            
            # Add image - fit to slide while preserving aspect ratio
            left = Inches(1)
            top = Inches(1.5)
            max_width = Inches(8)
            max_height = Inches(5.5)
            
            # Add picture with only width specified first to get aspect ratio preserved
            pic = slide.shapes.add_picture(img_stream, left, top, width=max_width)
            
            # Check if it exceeds max height
            if pic.height > max_height:
                # Resize by height
                ratio = max_height / pic.height
                pic.height = max_height
                pic.width = int(pic.width * ratio)
            
            # Center the image horizontally
            slide_width = prs.slide_width
            pic.left = int((slide_width - pic.width) / 2)
            
            # Add notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = "Proposed High-Level Architecture Diagram."
            
        except Exception as e:
            logger.error(f"Failed to add diagram slide: {e}")

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
